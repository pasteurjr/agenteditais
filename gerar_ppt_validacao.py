#!/usr/bin/env python3
"""
Gera o PowerPoint comercial: Facilitia.ia — Sprints 1 a 5
Saída: docs/APRESENTACAO_VALIDACAO_SPRINTS_1A5.pptx

Versão visual: fundo claro, verde esmeralda, cards com ícones,
callouts nos screenshots, diagramas nativos, excerpts dos docs V4.

Screenshots corretas por sprint:
  Sprint 1: runtime/screenshots/tutorial-sprint1-1/ (padrão F##_*)
  Sprint 2: runtime/screenshots/validacao-sprint2/  (padrão CV##-P##_*)
  Sprint 3: runtime/screenshots/validacao-sprint3/  (padrão P##_* e R##_*)
  Sprint 4: runtime/screenshots/validacao-sprint4/  (padrão I##_*, RE##_*, FU##_*)
  Sprint 5: runtime/screenshots/UC-CT##/ e UC-CRM##/ (padrão P##_resp.png)
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Constantes ──────────────────────────────────────────────

BASE = os.path.dirname(os.path.abspath(__file__))
SS_DIR = os.path.join(BASE, "runtime", "screenshots")
OUT = os.path.join(BASE, "docs", "APRESENTACAO_VALIDACAO_SPRINTS_1A5_V3.pptx")

# ── Paleta clara (comercial) ─────────────────────────────────
BG_LIGHT       = RGBColor(0xF8, 0xFA, 0xFC)  # slate-50
BG_CARD        = RGBColor(0xFF, 0xFF, 0xFF)  # white
BG_CARD_ALT    = RGBColor(0xF1, 0xF5, 0xF9)  # slate-100
BORDER_SOFT    = RGBColor(0xE2, 0xE8, 0xF0)  # slate-200

# Primária: verde esmeralda
EMERALD        = RGBColor(0x05, 0x96, 0x69)  # emerald-600 — primária
EMERALD_DARK   = RGBColor(0x04, 0x7A, 0x55)  # emerald-700
EMERALD_LIGHT  = RGBColor(0xD1, 0xFA, 0xE5)  # emerald-100

# Texto
TXT_INK        = RGBColor(0x0F, 0x17, 0x2A)  # slate-900 — títulos
TXT_BODY       = RGBColor(0x33, 0x41, 0x55)  # slate-700 — texto
TXT_MUTED      = RGBColor(0x64, 0x74, 0x8B)  # slate-500 — rodapés

# Acentos
ACCENT_RED     = RGBColor(0xDC, 0x26, 0x26)  # para callouts
ACCENT_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_BLUE    = RGBColor(0x25, 0x63, 0xEB)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

# Cores por sprint (tons de acento)
SPRINT_COLORS = {
    1: RGBColor(0x25, 0x63, 0xEB),  # blue-600
    2: RGBColor(0x05, 0x96, 0x69),  # emerald-600
    3: RGBColor(0x7C, 0x3A, 0xED),  # violet-600
    4: RGBColor(0xDC, 0x26, 0x26),  # red-600
    5: RGBColor(0xEA, 0x58, 0x0C),  # orange-600
}

SPRINT_ICONS = {
    1: "🏗️",
    2: "🔍",
    3: "💰",
    4: "⚖️",
    5: "🤝",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Mapeamento de screenshots (preservado da versão anterior) ──

SPRINT1_SS = {
    "UC-F01": "tutorial-sprint1-1/F01_02_razao_social.png",
    "UC-F02": "tutorial-sprint1-1/F02_04_area_padrao.png",
    "UC-F03": "tutorial-sprint1-1/F03_05_lista_documentos.png",
    "UC-F04": "tutorial-sprint1-1/F04_04_busca_concluida.png",
    "UC-F05": "tutorial-sprint1-1/F05_05_lista_responsaveis.png",
    "UC-F06": "VC/sprint1/F06_detalhe_produto.png",
    "UC-F07": "VC/sprint1/F07_com_produtos.png",
    "UC-F08": "tutorial-sprint1-1/F08_04_produto_salvo.png",
    "UC-F09": "tutorial-sprint1-1/F09_03_reprocessado.png",
    "UC-F10": "tutorial-sprint1-1/F10_04_resultado_web.png",
    "UC-F11": "tutorial-sprint1-1/F11_02_scores.png",
    "UC-F12": "VC/sprint1/F12_detalhes_metadados.png",
    "UC-F13": "tutorial-sprint1-1/F13_04_funil_monitoramento.png",
    "UC-F14": "tutorial-sprint1-1/F14_06_limiares_salvos.png",
    "UC-F15": "tutorial-sprint1-1/F15_07_modalidades_salvas.png",
    "UC-F16": "tutorial-sprint1-1/F16_07_ncms_salvos.png",
    "UC-F17": "tutorial-sprint1-1/F17_02_notificacoes_preenchidas.png",
}

SPRINT2_SS = {
    "UC-CV01": "validacao-sprint2/CV01-P05_resp_resultados.png",
    "UC-CV02": "validacao-sprint2/CV02_resp_painel.png",
    "UC-CV03": "validacao-sprint2/CV03-P01_resp_salvo.png",
    "UC-CV04": "validacao-sprint2/CV04_resp_salva.png",
    "UC-CV05": "validacao-sprint2/CV05_resp_exportado.png",
    "UC-CV06": "validacao-sprint2/CV06_resp_criado.png",
    "UC-CV07": "validacao-sprint2/CV07_resp_selecionado.png",
    "UC-CV08": "validacao-sprint2/CV08_resp_scores.png",
    "UC-CV09": "validacao-sprint2/CV09_resp_itens.png",
    "UC-CV10": "validacao-sprint2/CV10_resp_documentacao.png",
    "UC-CV11": "validacao-sprint2/CV11_resp_vencedores.png",
    "UC-CV12": "validacao-sprint2/CV12_resp_mercado.png",
    "UC-CV13": "validacao-sprint2/CV13_resp_resposta.png",
}

SPRINT3_SS = {
    "UC-P01": "validacao-sprint3/P01_p7_tabela_itens.png",
    "UC-P02": "validacao-sprint3/P02_resp_vinculacao.png",
    "UC-P03": "validacao-sprint3/P03_p6_volumetria_calculada.png",
    "UC-P04": "validacao-sprint3/P04_p6_custos_salvos.png",
    "UC-P05": "validacao-sprint3/P05_resp_preco_base_salvo.png",
    "UC-P06": "validacao-sprint3/P06_resp_referencia_salva.png",
    "UC-P07": "validacao-sprint3/P07_resp_lances_salvos.png",
    "UC-P08": "validacao-sprint3/P08_p9_simulador_disputa.png",
    "UC-P09": "validacao-sprint3/P09_p3_resultado_filtro.png",
    "UC-P10": "validacao-sprint3/P10_p4_comodato_salvo.png",
    "UC-P11": "validacao-sprint3/P11_p7_insights_regenerados.png",
    "UC-P12": "validacao-sprint3/P12_p3_relatorio_html.png",
    "UC-R01": "validacao-sprint3/R01_p9_proposta_gerada.png",
    "UC-R02": "validacao-sprint3/R02_p5_importado.png",
    "UC-R03": "validacao-sprint3/R03_p5_conteudo_salvo.png",
    "UC-R04": "validacao-sprint3/R04_p3_tabela_anvisa.png",
    "UC-R05": "validacao-sprint3/R05_p4_tabela_documentos.png",
    "UC-R06": "validacao-sprint3/R06_p1_card_exportacao.png",
    "UC-R07": "validacao-sprint3/R07_resp_marcada_enviada.png",
}

SPRINT4_SS = {
    "UC-I01": "validacao-sprint4/I01_resp_tabela_inconsistencias.png",
    "UC-I02": "validacao-sprint4/I02_resp_checkboxes_marcados.png",
    "UC-I03": "validacao-sprint4/I03_resp_peticao_criada.png",
    "UC-I04": "validacao-sprint4/I04_resp_upload.png",
    "UC-I05": "validacao-sprint4/I05_resp_countdown.png",
    "UC-RE01": "validacao-sprint4/RE01_resp_monitoramento_ativado.png",
    "UC-RE02": "validacao-sprint4/RE02_resp_analise_vencedora.png",
    "UC-RE03": "validacao-sprint4/RE03_resp_resposta_chat.png",
    "UC-RE04": "validacao-sprint4/RE04_resp_laudo_criado.png",
    "UC-RE05": "validacao-sprint4/RE05_resp_contra_razao_criada.png",
    "UC-RE06": "validacao-sprint4/RE06_resp_submissao_registrada.png",
    "UC-FU01": "validacao-sprint4/FU01_resp_resultado_ganho.png",
    "UC-FU02": "validacao-sprint4/FU02_resp_alertas_configurados.png",
    "UC-FU03": "validacao-sprint4/FU03_resp_score_logistico.png",
}

SCREENSHOT_OVERRIDES = {
    "UC-CR01-acao": "UC-CR01/P01_acao.png",
    "UC-CRM03": "UC-CRM03/P02_resp.png",
    "UC-CRM04": "UC-CRM04/P03_resp.png",
}


def screenshot_path(uc_code):
    """Retorna o caminho correto da screenshot para o UC."""
    if uc_code in SCREENSHOT_OVERRIDES:
        p = os.path.join(SS_DIR, SCREENSHOT_OVERRIDES[uc_code])
        return p if os.path.exists(p) else None
    if uc_code in SPRINT1_SS:
        p = os.path.join(SS_DIR, SPRINT1_SS[uc_code])
        return p if os.path.exists(p) else None
    if uc_code in SPRINT2_SS:
        p = os.path.join(SS_DIR, SPRINT2_SS[uc_code])
        return p if os.path.exists(p) else None
    if uc_code in SPRINT3_SS:
        p = os.path.join(SS_DIR, SPRINT3_SS[uc_code])
        return p if os.path.exists(p) else None
    if uc_code in SPRINT4_SS:
        p = os.path.join(SS_DIR, SPRINT4_SS[uc_code])
        return p if os.path.exists(p) else None
    d = os.path.join(SS_DIR, uc_code)
    if os.path.isdir(d):
        for pat in ["P03_resp*", "P02_resp*", "P01_resp*", "*_resp*"]:
            matches = sorted(glob.glob(os.path.join(d, pat)))
            if matches:
                return matches[0]
        pngs = sorted(glob.glob(os.path.join(d, "*.png")))
        return pngs[0] if pngs else None
    return None


# ═══════════════════════════════════════════════════════════
# HELPERS VISUAIS
# ═══════════════════════════════════════════════════════════

def set_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=14,
                color=TXT_BODY, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    return tb


def add_multiline(slide, left, top, width, height, lines, *, size=14,
                  color=TXT_BODY, bold=False, align=PP_ALIGN.LEFT,
                  line_spacing=1.15, font="Inter"):
    """lines: lista de (texto, bold?, size?, color?) ou string simples."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, sz, cl = item, bold, size, color
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            sz = item[2] if len(item) > 2 else size
            cl = item[3] if len(item) > 3 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.alignment = align
        p.line_spacing = line_spacing
        p.font.size = Pt(sz)
        p.font.color.rgb = cl
        p.font.bold = b
        p.font.name = font
    return tb


def add_rect(slide, left, top, width, height, *, fill=BG_CARD,
             border=BORDER_SOFT, border_w=0.75, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    if not shadow:
        sp = shp.shadow
        try:
            sp.inherit = False
        except Exception:
            pass
    return shp


def add_band(slide, left, top, width, height, color):
    """Retângulo sólido sem borda (faixa/acento)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_pill(slide, left, top, width, height, text, *, fill=EMERALD,
             text_color=WHITE, size=11, bold=True):
    """Badge arredondada com texto centralizado."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = "Inter"
    return shp


def add_circle_num(slide, left, top, diameter, number, *, fill=ACCENT_RED, text_color=WHITE):
    """Círculo numerado — para callouts."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = WHITE
    shp.line.width = Pt(2)
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Inter"
    return shp


def add_highlight_box(slide, left, top, width, height, *, color=ACCENT_RED, border_w=2.5):
    """Retângulo vazio com borda colorida — destaca área de screenshot."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.05
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(border_w)
    return shp


def add_screenshot(slide, img_path, left, top, max_w, max_h):
    """Insere screenshot com moldura branca e sombra leve."""
    if not img_path or not os.path.exists(img_path):
        # Placeholder se faltar
        add_rect(slide, left, top, max_w, max_h, fill=BG_CARD_ALT, border=BORDER_SOFT)
        add_textbox(slide, left, top + max_h/2 - Inches(0.15), max_w, Inches(0.3),
                    "screenshot indisponível", size=11, color=TXT_MUTED,
                    align=PP_ALIGN.CENTER)
        return None
    # Moldura branca
    frame_pad = Inches(0.08)
    add_rect(slide, left - frame_pad, top - frame_pad,
             max_w + 2*frame_pad, max_h + 2*frame_pad,
             fill=BG_CARD, border=BORDER_SOFT, border_w=0.5)
    try:
        pic = slide.shapes.add_picture(img_path, left, top, width=max_w, height=max_h)
        return pic
    except Exception as e:
        print(f"⚠️  Erro inserindo {img_path}: {e}")
        return None


def add_footer(slide, text, *, sprint_color=None):
    """Rodapé fino com linha colorida."""
    # Linha
    line_color = sprint_color or EMERALD
    add_band(slide, Inches(0.5), Inches(7.15), Inches(12.33), Inches(0.02), line_color)
    # Texto
    add_textbox(slide, Inches(0.5), Inches(7.22), Inches(11.5), Inches(0.25),
                text, size=9, color=TXT_MUTED)
    # Logo
    add_textbox(slide, Inches(11.8), Inches(7.22), Inches(1.5), Inches(0.25),
                "Facilitia.ia", size=9, color=EMERALD, bold=True,
                align=PP_ALIGN.RIGHT)


def slide_header(slide, kicker, title, subtitle=None, *, sprint_num=None):
    """Cabeçalho com kicker colorido + título + subtítulo."""
    color = SPRINT_COLORS[sprint_num] if sprint_num else EMERALD
    icon = SPRINT_ICONS[sprint_num] if sprint_num else "🎯"

    # Kicker (badge pequena no topo)
    kicker_w = Inches(max(1.4, 0.13 * len(kicker) + 0.5))
    add_pill(slide, Inches(0.6), Inches(0.35), kicker_w, Inches(0.32),
             f"{icon}  {kicker}", fill=color, text_color=WHITE, size=10, bold=True)

    # Título principal — tamanho reduzido e mais altura pra não encostar no subtítulo
    add_textbox(slide, Inches(0.6), Inches(0.78), Inches(12.2), Inches(0.85),
                title, size=26, color=TXT_INK, bold=True)

    # Subtítulo
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.65), Inches(12.2), Inches(0.35),
                    subtitle, size=13, color=TXT_MUTED)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    return slide


# ═══════════════════════════════════════════════════════════
# TEMPLATES DE SLIDE
# ═══════════════════════════════════════════════════════════

def cover_slide(prs):
    """Capa principal — fundo com acento verde."""
    slide = new_slide(prs)

    # Faixa vertical esquerda
    add_band(slide, Inches(0), Inches(0), Inches(0.3), SLIDE_H, EMERALD)

    # Logo/brand topo
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
                "FACILITIA.IA", size=12, color=EMERALD, bold=True)

    # Título gigante
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(1.2),
                "Apresentação de Validação", size=40, color=TXT_INK, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.55), Inches(12), Inches(1.0),
                "Sprints 1 a 5", size=56, color=EMERALD, bold=True)

    # Barra decorativa
    add_band(slide, Inches(0.8), Inches(3.85), Inches(2.0), Inches(0.06), EMERALD)

    # Subtítulo
    add_textbox(slide, Inches(0.8), Inches(4.05), Inches(12), Inches(0.5),
                "Sistema completo de gestão de licitações com IA",
                size=18, color=TXT_BODY)

    # Métricas em pills horizontais
    metrics = [
        ("79", "Casos de uso"),
        ("100%", "Aprovação"),
        ("217", "Regras de negócio"),
        ("322", "Testes Playwright"),
    ]
    x = Inches(0.8)
    for num, label in metrics:
        add_rect(slide, x, Inches(4.9), Inches(2.7), Inches(1.1),
                 fill=BG_CARD, border=BORDER_SOFT)
        add_textbox(slide, x, Inches(5.05), Inches(2.7), Inches(0.5),
                    num, size=28, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.55), Inches(2.7), Inches(0.35),
                    label, size=11, color=TXT_MUTED, align=PP_ALIGN.CENTER)
        x += Inches(2.9)

    # Rodapé
    add_textbox(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
                "Abril 2026  •  Documento de validação comercial",
                size=12, color=TXT_MUTED)

    return slide


def sprint_cover(prs, sprint_num, title, subtitle, bullets):
    """Capa de sprint — fundo branco, acento colorido vertical."""
    slide = new_slide(prs)
    color = SPRINT_COLORS[sprint_num]
    icon = SPRINT_ICONS[sprint_num]

    # Faixa vertical colorida
    add_band(slide, Inches(0), Inches(0), Inches(3.0), SLIDE_H, color)

    # Número gigante do sprint (fundo)
    add_textbox(slide, Inches(0.3), Inches(1.4), Inches(2.5), Inches(3.5),
                str(sprint_num), size=180, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.3), Inches(5.2), Inches(2.5), Inches(0.5),
                "SPRINT", size=20, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.3), Inches(5.7), Inches(2.5), Inches(0.5),
                icon, size=30, color=WHITE, align=PP_ALIGN.CENTER)

    # Título + subtítulo à direita — altura maior pra título quebrar em 2 linhas
    add_textbox(slide, Inches(3.5), Inches(1.3), Inches(9.6), Inches(1.7),
                title, size=30, color=TXT_INK, bold=True)
    add_textbox(slide, Inches(3.5), Inches(2.95), Inches(9.6), Inches(0.5),
                subtitle, size=14, color=color, bold=True)

    # Linha divisora
    add_band(slide, Inches(3.5), Inches(3.55), Inches(1.5), Inches(0.05), color)

    # Bullets como cards
    y = Inches(3.85)
    for b in bullets:
        add_textbox(slide, Inches(3.5), y, Inches(0.3), Inches(0.4),
                    "▸", size=18, color=color, bold=True)
        add_textbox(slide, Inches(3.9), y, Inches(9.1), Inches(0.4),
                    b, size=14, color=TXT_BODY)
        y += Inches(0.5)

    add_footer(slide, f"Sprint {sprint_num}  •  Visão geral", sprint_color=color)
    return slide


def content_slide(prs, *, kicker, title, subtitle, bullets,
                  screenshot_uc=None, callouts=None, footer=None,
                  sprint_num=None, excerpt=None):
    """
    Slide padrão com bullets à esquerda e screenshot à direita.
    callouts: lista de (x_rel, y_rel, w_rel, h_rel, num, explicacao)
              coordenadas relativas (0-1) à área do screenshot.
    excerpt: dict {"title": "...", "lines": [...], "source": "..."} para caixa de trecho.
    """
    slide = new_slide(prs)
    slide_header(slide, kicker, title, subtitle, sprint_num=sprint_num)

    has_screenshot = bool(screenshot_uc)
    img_path = screenshot_path(screenshot_uc) if screenshot_uc else None

    if has_screenshot:
        # Layout 2 colunas: bullets 5.2in | screenshot 7.4in
        bullets_left   = Inches(0.6)
        bullets_width  = Inches(5.2)
        shot_left      = Inches(6.0)
        shot_top       = Inches(2.1)
        shot_width     = Inches(6.8)
        shot_height    = Inches(4.7)

        # Card de bullets
        add_rect(slide, bullets_left, Inches(2.05), bullets_width, Inches(4.9),
                 fill=BG_CARD, border=BORDER_SOFT)
        # Bullets com bolinhas coloridas
        y = Inches(2.3)
        color = SPRINT_COLORS[sprint_num] if sprint_num else EMERALD
        for b in bullets:
            # bullet point
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bullets_left + Inches(0.25),
                                          y + Inches(0.08), Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            # texto
            add_textbox(slide, bullets_left + Inches(0.5), y,
                        bullets_width - Inches(0.6), Inches(0.9),
                        b, size=12, color=TXT_BODY)
            y += Inches(0.75)

        # Screenshot
        add_screenshot(slide, img_path, shot_left, shot_top, shot_width, shot_height)

        # Callouts sobre o screenshot
        if callouts and img_path:
            for i, co in enumerate(callouts, start=1):
                x_rel, y_rel, w_rel, h_rel, explicacao = co
                hl_left = shot_left + Emu(int(shot_width * x_rel))
                hl_top  = shot_top + Emu(int(shot_height * y_rel))
                hl_w    = Emu(int(shot_width * w_rel))
                hl_h    = Emu(int(shot_height * h_rel))
                add_highlight_box(slide, hl_left, hl_top, hl_w, hl_h)
                # círculo numerado no canto sup-esq do highlight
                add_circle_num(slide, hl_left - Inches(0.15), hl_top - Inches(0.15),
                               Inches(0.35), i)

            # Legenda abaixo do screenshot — explicações numeradas
            lgd_top = shot_top + shot_height + Inches(0.15)
            if lgd_top + Inches(0.5) <= Inches(7.0):
                txt = "   ".join([f"❶❷❸❹❺❻❼❽❾"[i-1] + " " + co[4]
                                  for i, co in enumerate(callouts, start=1)
                                  if i <= 9])
                # Fallback ASCII se emojis de círculo não renderizarem
                parts = []
                for i, co in enumerate(callouts, start=1):
                    parts.append(f"({i}) {co[4]}")
                add_textbox(slide, shot_left, lgd_top, shot_width, Inches(0.35),
                            "   ".join(parts), size=9, color=TXT_MUTED)

    else:
        # Só bullets, coluna única larga
        bullets_left  = Inches(0.6)
        bullets_width = Inches(12.2)
        add_rect(slide, bullets_left, Inches(2.05), bullets_width, Inches(4.9),
                 fill=BG_CARD, border=BORDER_SOFT)
        color = SPRINT_COLORS[sprint_num] if sprint_num else EMERALD
        y = Inches(2.4)
        for b in bullets:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bullets_left + Inches(0.5),
                                          y + Inches(0.12), Inches(0.15), Inches(0.15))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            add_textbox(slide, bullets_left + Inches(0.85), y,
                        bullets_width - Inches(1.2), Inches(0.8),
                        b, size=15, color=TXT_BODY)
            y += Inches(0.85)

    # Caixa de trecho do documento (excerpt)
    if excerpt:
        exc_left = Inches(0.6)
        exc_top  = Inches(6.0) if not has_screenshot else Inches(5.95)
        exc_w    = Inches(5.2) if has_screenshot else Inches(12.2)
        exc_h    = Inches(1.05)
        add_rect(slide, exc_left, exc_top, exc_w, exc_h,
                 fill=EMERALD_LIGHT, border=EMERALD, border_w=1)
        add_textbox(slide, exc_left + Inches(0.15), exc_top + Inches(0.08),
                    exc_w - Inches(0.3), Inches(0.3),
                    f"📄 {excerpt['title']}", size=10, color=EMERALD_DARK, bold=True)
        body_txt = " / ".join(excerpt.get("lines", []))[:200]
        add_textbox(slide, exc_left + Inches(0.15), exc_top + Inches(0.35),
                    exc_w - Inches(0.3), Inches(0.55),
                    body_txt, size=9, color=TXT_BODY, font="Courier New")
        if excerpt.get("source"):
            add_textbox(slide, exc_left + Inches(0.15), exc_top + Inches(0.82),
                        exc_w - Inches(0.3), Inches(0.2),
                        f"Fonte: {excerpt['source']}", size=8, color=TXT_MUTED)

    add_footer(slide, footer or "", sprint_color=SPRINT_COLORS.get(sprint_num))
    return slide


def diagram_slide(prs, *, kicker, title, subtitle, draw_fn, footer=None, sprint_num=None):
    """Slide com área livre para desenhar diagrama nativo."""
    slide = new_slide(prs)
    slide_header(slide, kicker, title, subtitle, sprint_num=sprint_num)
    draw_fn(slide)
    add_footer(slide, footer or "", sprint_color=SPRINT_COLORS.get(sprint_num))
    return slide


# ═══════════════════════════════════════════════════════════
# DIAGRAMAS NATIVOS
# ═══════════════════════════════════════════════════════════

def draw_jornada_editais(slide):
    """Funil: Descoberta → Análise → Precificação → Proposta → Execução."""
    etapas = [
        ("1", "Descoberta",     "Busca PNCP + alerts",              SPRINT_COLORS[2], "🔍"),
        ("2", "Análise/GO",     "Score 6D + decisão",               SPRINT_COLORS[2], "🎯"),
        ("3", "Precificação",   "6 camadas até lance",              SPRINT_COLORS[3], "💰"),
        ("4", "Proposta",       "Dossiê técnico",                   SPRINT_COLORS[3], "📋"),
        ("5", "Submissão",      "Portal gov.br",                    SPRINT_COLORS[4], "📤"),
        ("6", "Execução/CRM",   "Contrato + pós-venda",             SPRINT_COLORS[5], "🤝"),
    ]
    top = Inches(2.4)
    card_w = Inches(2.0)
    card_h = Inches(3.3)
    gap = Inches(0.1)
    start_x = Inches(0.55)
    for i, (num, nome, desc, color, icon) in enumerate(etapas):
        x = start_x + (card_w + gap) * i
        # Card
        add_rect(slide, x, top, card_w, card_h, fill=BG_CARD, border=BORDER_SOFT)
        # Faixa colorida no topo
        add_band(slide, x, top, card_w, Inches(0.08), color)
        # Número grande
        add_textbox(slide, x, top + Inches(0.25), card_w, Inches(0.7),
                    num, size=44, color=color, bold=True, align=PP_ALIGN.CENTER)
        # Ícone
        add_textbox(slide, x, top + Inches(1.0), card_w, Inches(0.6),
                    icon, size=32, align=PP_ALIGN.CENTER)
        # Nome
        add_textbox(slide, x, top + Inches(1.7), card_w, Inches(0.4),
                    nome, size=14, color=TXT_INK, bold=True, align=PP_ALIGN.CENTER)
        # Descrição
        add_textbox(slide, x + Inches(0.1), top + Inches(2.15), card_w - Inches(0.2), Inches(0.9),
                    desc, size=10, color=TXT_MUTED, align=PP_ALIGN.CENTER)
        # Seta entre cards
        if i < len(etapas) - 1:
            arrow_x = x + card_w - Inches(0.02)
            add_textbox(slide, arrow_x, top + Inches(1.1), gap + Inches(0.2), Inches(0.5),
                        "▶", size=16, color=TXT_MUTED, align=PP_ALIGN.CENTER)


def draw_arquitetura(slide):
    """3 camadas: Frontend → Backend → Serviços externos."""
    top = Inches(2.3)
    card_h = Inches(1.35)
    gap = Inches(0.25)
    camadas = [
        ("🖥️  Frontend",    "React + Vite  •  porta 5180",          "ProducaoPage, CRMPage, RecursosPage, ValidacaoPage, PrecificacaoPage, PropostaPage",     ACCENT_BLUE),
        ("⚙️  Backend",     "Flask + SQLAlchemy  •  porta 5007",    "backend/app.py, crud_routes.py, empenho_routes.py, crm_routes.py, tools.py",             EMERALD),
        ("🌐  Externos",    "DeepSeek IA  •  PNCP  •  Brave Search",  "OpenRouter LLM, ANVISA, RF/gov.br, Stripe, SMTP",                                       RGBColor(0x7C, 0x3A, 0xED)),
    ]
    for i, (titulo, sub, itens, color) in enumerate(camadas):
        y = top + (card_h + gap) * i
        # Card
        add_rect(slide, Inches(0.6), y, Inches(12.1), card_h,
                 fill=BG_CARD, border=BORDER_SOFT)
        # Faixa colorida lateral
        add_band(slide, Inches(0.6), y, Inches(0.15), card_h, color)
        # Título
        add_textbox(slide, Inches(0.95), y + Inches(0.18), Inches(4.5), Inches(0.4),
                    titulo, size=16, color=TXT_INK, bold=True)
        # Sub
        add_textbox(slide, Inches(0.95), y + Inches(0.58), Inches(4.5), Inches(0.3),
                    sub, size=11, color=color, bold=True)
        # Itens
        add_textbox(slide, Inches(5.5), y + Inches(0.3), Inches(7.1), Inches(0.8),
                    itens, size=11, color=TXT_BODY)
        # Seta pra baixo
        if i < len(camadas) - 1:
            add_textbox(slide, Inches(6.5), y + card_h - Inches(0.05),
                        Inches(0.3), gap, "▼", size=14, color=TXT_MUTED, align=PP_ALIGN.CENTER)


def draw_pipeline_validacao(slide):
    """Ciclo: Preparador → Especificador → Executor → Relator + Humano."""
    nodes = [
        ("🧠", "Preparador\nChatGPT",   "Lê docs, gera\nanálise estrutural",      ACCENT_BLUE),
        ("📝", "Especificador\nChatGPT", "Cria spec de\ntestes + dados",          RGBColor(0x7C, 0x3A, 0xED)),
        ("🤖", "Executor\nClaude",      "Roda Playwright,\ncaptura screenshots", EMERALD),
        ("📊", "Relator\nClaude",       "Compila resultado,\ngera PPT + PDFs",   ACCENT_AMBER),
    ]
    top = Inches(2.6)
    node_w = Inches(2.55)
    node_h = Inches(2.5)
    gap = Inches(0.6)
    start_x = Inches(0.9)
    for i, (icon, titulo, desc, color) in enumerate(nodes):
        x = start_x + (node_w + gap) * i
        # Circle background
        add_rect(slide, x, top, node_w, node_h, fill=BG_CARD, border=color, border_w=2)
        # Ícone
        add_textbox(slide, x, top + Inches(0.25), node_w, Inches(0.8),
                    icon, size=44, align=PP_ALIGN.CENTER)
        # Título
        add_textbox(slide, x, top + Inches(1.0), node_w, Inches(0.7),
                    titulo, size=14, color=TXT_INK, bold=True, align=PP_ALIGN.CENTER)
        # Descrição
        add_textbox(slide, x, top + Inches(1.65), node_w, Inches(0.8),
                    desc, size=10, color=TXT_MUTED, align=PP_ALIGN.CENTER)
        # Seta
        if i < len(nodes) - 1:
            add_textbox(slide, x + node_w, top + Inches(1.0), gap, Inches(0.6),
                        "→", size=28, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)

    # Validador humano (bloco final embaixo)
    add_rect(slide, Inches(0.9), Inches(5.5), Inches(12.0), Inches(1.3),
             fill=EMERALD_LIGHT, border=EMERALD, border_w=1.5)
    add_textbox(slide, Inches(1.1), Inches(5.65), Inches(1.0), Inches(1.0),
                "👤", size=44, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.2), Inches(5.7), Inches(10.5), Inches(0.4),
                "Validador humano (valida2@valida.com.br)",
                size=15, color=EMERALD_DARK, bold=True)
    add_textbox(slide, Inches(2.2), Inches(6.1), Inches(10.5), Inches(0.8),
                "Segue tutorialsprint#-2.md como roteiro com dados da empresa RP3X — reporta divergências no spreadsheet de validação.",
                size=11, color=TXT_BODY)


def draw_torre_precificacao(slide):
    """6 camadas de A a F — torre empilhada (do topo F para a base A)."""
    camadas = [
        ("F", "Histórico",           "UC-P09 • Preços pagos anteriores",  RGBColor(0x7C, 0x3A, 0xED)),
        ("E", "Lance Mínimo",        "UC-P08 • Piso de disputa",          RGBColor(0x67, 0x2B, 0xE0)),
        ("D", "Lance Inicial",       "UC-P07 • Abertura do pregão",       RGBColor(0x52, 0x24, 0xD0)),
        ("C", "Valor de Referência", "UC-P06 • Preço defensável",         RGBColor(0x3E, 0x1C, 0xB8)),
        ("B", "Preço Base",          "UC-P05 • Custo + margem",           RGBColor(0x2D, 0x14, 0x9E)),
        ("A", "Volumetria",          "UC-P03 • Quantidade técnica",       RGBColor(0x1E, 0x0C, 0x7D)),
    ]
    top = Inches(2.2)
    card_h = Inches(0.75)
    gap = Inches(0.04)
    for i, (letra, nome, desc, color) in enumerate(camadas):
        y = top + (card_h + gap) * i
        # largura decrescente para baixo (topo mais largo, base menor)? Não - visualmente, Camada A é base (maior)
        # Vamos inverter: A no topo visual (topo da lista), base visual. Aqui 0 é F (topo narrativo).
        add_rect(slide, Inches(3.5), y, Inches(6.3), card_h, fill=color, border=None)
        add_textbox(slide, Inches(3.7), y + Inches(0.15), Inches(0.8), Inches(0.45),
                    letra, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.5), y + Inches(0.1), Inches(2.8), Inches(0.35),
                    nome, size=14, color=WHITE, bold=True)
        add_textbox(slide, Inches(4.5), y + Inches(0.4), Inches(5.2), Inches(0.3),
                    desc, size=10, color=WHITE)

    # Seta lateral "fluxo"
    add_textbox(slide, Inches(10.0), Inches(2.4), Inches(0.5), Inches(4.8),
                "↓", size=40, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.5), Inches(2.5), Inches(2.5), Inches(0.4),
                "Fluxo do cálculo", size=11, color=EMERALD, bold=True)
    add_textbox(slide, Inches(10.5), Inches(2.9), Inches(2.5), Inches(3.0),
                "Cada camada usa dados das anteriores. A Camada A parte da volumetria técnica; cada passo adiciona contexto até produzir o lance mínimo executável.",
                size=10, color=TXT_BODY)


def draw_anatomia_uc(slide):
    """Anatomia de um caso de uso — com UC-CV08 como exemplo."""
    # Card principal
    add_rect(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.85),
             fill=BG_CARD, border=BORDER_SOFT)

    # Cabeçalho do UC
    add_band(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(0.5), EMERALD)
    add_textbox(slide, Inches(0.85), Inches(2.18), Inches(11.5), Inches(0.35),
                "## [UC-CV08] Calcular scores multidimensionais e decidir GO/NO-GO",
                size=14, color=WHITE, bold=True, font="Courier New")

    # Conteúdo em 2 colunas
    secoes_esq = [
        ("**RNs aplicadas:**", "RN-047, RN-048, RN-049, RN-050, RN-051, RN-052, RN-053, RN-054"),
        ("**RF relacionados:**", "RF-027 (score multidimensional), RF-028 (decisão GO/NO-GO), RF-037 (auditoria de decisão)"),
        ("**Ator:**", "Usuário analista/comercial"),
    ]
    y = Inches(2.8)
    for label, val in secoes_esq:
        add_textbox(slide, Inches(0.85), y, Inches(5.8), Inches(0.3),
                    label.replace("**", ""), size=11, color=EMERALD_DARK, bold=True,
                    font="Courier New")
        add_textbox(slide, Inches(0.85), y + Inches(0.3), Inches(5.8), Inches(0.5),
                    val, size=10, color=TXT_BODY)
        y += Inches(0.85)

    # Seção direita: Sequência de eventos
    add_textbox(slide, Inches(6.9), Inches(2.8), Inches(5.7), Inches(0.3),
                "### Sequência de eventos", size=11, color=EMERALD_DARK, bold=True,
                font="Courier New")
    passos = [
        "1. Abrir painel de Validação",
        "2. Selecionar edital para análise",
        "3. Sistema calcula 6 dimensões (técnica, documental, financeira, experiência, geográfica, prazo)",
        "4. Exibir scores 0-100 por dimensão",
        "5. Sistema sugere GO ou NO-GO",
        "6. Usuário confirma ou reverte decisão",
        "7. Sistema registra justificativa",
        "8. Auditoria grava decisão + RNs aplicadas",
    ]
    y = Inches(3.15)
    for p in passos:
        add_textbox(slide, Inches(6.9), y, Inches(5.7), Inches(0.25),
                    p, size=9, color=TXT_BODY)
        y += Inches(0.28)

    # Caixa de legenda embaixo
    add_rect(slide, Inches(0.6), Inches(6.08), Inches(12.1), Inches(0.85),
             fill=EMERALD_LIGHT, border=EMERALD, border_w=1)
    add_textbox(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.3),
                "💡  Todo UC segue este mesmo padrão: RNs, RF, Ator, Pré/Pós-condições, Sequência, Tela, Implementação atual.",
                size=11, color=EMERALD_DARK, bold=True)
    add_textbox(slide, Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.4),
                "Fonte: docs/CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V4.md, linha 731. Mesmo schema em todas as 5 sprints (V4).",
                size=9, color=TXT_BODY)


def draw_pipeline_crm(slide):
    """13 stages do CRM em grid horizontal."""
    stages = [
        ("Captado",     14),
        ("Análise",     22),
        ("GO",          18),
        ("NO-GO",       9),
        ("Preço",       11),
        ("Proposta",    13),
        ("Submetido",   8),
        ("Disputa",     6),
        ("Impugn.",     3),
        ("Recurso",     2),
        ("Vencido",     4),
        ("Perdido",     3),
        ("Concluído",   8),
    ]
    top = Inches(2.4)
    card_w = Inches(0.92)
    card_h = Inches(2.3)
    gap = Inches(0.05)
    start_x = Inches(0.55)
    for i, (nome, cnt) in enumerate(stages):
        x = start_x + (card_w + gap) * i
        color = SPRINT_COLORS[5] if i < 4 else (EMERALD if i < 8 else ACCENT_AMBER)
        add_rect(slide, x, top, card_w, card_h, fill=BG_CARD, border=BORDER_SOFT)
        add_band(slide, x, top, card_w, Inches(0.12), color)
        # Número do stage
        add_textbox(slide, x, top + Inches(0.22), card_w, Inches(0.35),
                    f"#{i+1}", size=10, color=TXT_MUTED, bold=True, align=PP_ALIGN.CENTER)
        # Nome
        add_textbox(slide, x + Inches(0.03), top + Inches(0.6), card_w - Inches(0.06), Inches(0.5),
                    nome, size=11, color=TXT_INK, bold=True, align=PP_ALIGN.CENTER)
        # Contador
        add_textbox(slide, x, top + Inches(1.25), card_w, Inches(0.55),
                    str(cnt), size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(1.85), card_w, Inches(0.3),
                    "editais", size=8, color=TXT_MUTED, align=PP_ALIGN.CENTER)

    # Legenda
    add_textbox(slide, Inches(0.55), Inches(5.0), Inches(12.2), Inches(0.35),
                "Total: 151 editais distribuídos no pipeline (dados reais do seed CH Hospitalar)",
                size=12, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)

    # Blocos explicativos
    legendas = [
        ("🟧 Captação", "Stages 1-4: busca + triagem GO/NO-GO", SPRINT_COLORS[5]),
        ("🟢 Fluxo ativo", "Stages 5-8: precificação, proposta, disputa", EMERALD),
        ("🟡 Desfechos", "Stages 9-13: impugnação, recurso, resultado", ACCENT_AMBER),
    ]
    x = Inches(0.55)
    for titulo, desc, color in legendas:
        add_rect(slide, x, Inches(5.5), Inches(4.2), Inches(1.2),
                 fill=BG_CARD, border=color, border_w=2)
        add_textbox(slide, x + Inches(0.15), Inches(5.6), Inches(4.0), Inches(0.35),
                    titulo, size=13, color=color, bold=True)
        add_textbox(slide, x + Inches(0.15), Inches(5.95), Inches(4.0), Inches(0.7),
                    desc, size=10, color=TXT_BODY)
        x += Inches(4.3)


def draw_ciclo_testes(slide):
    """Ciclo completo: tutorial-1 → Playwright → screenshots → relatório."""
    nodes = [
        ("📖", "tutorial-1",       "Roteiro narrativo\nda CH Hospitalar",    ACCENT_BLUE, "Dados reais,\npasso a passo"),
        ("🤖", "Playwright",       "Automação\nsequencial",                  EMERALD, "login → ações\n→ asserts"),
        ("📸", "Screenshots",      "_acao + _resp\npor passo",                RGBColor(0x7C, 0x3A, 0xED), "Par A/R\nsempre"),
        ("✅", "Asserts",         "Banco + UI\n+ RN",                         ACCENT_AMBER, "9/9 verif\nDB PASS"),
        ("📄", "RESULTADO",        "VALIDACAO\nSPRINT#.md",                    ACCENT_RED, "Doc narrativa\nnum Aprovada"),
        ("👤", "Validador",        "tutorial-2\n(RP3X)",                      EMERALD_DARK, "Humano\nindependente"),
    ]
    top = Inches(2.3)
    node_w = Inches(2.0)
    node_h = Inches(2.1)
    gap = Inches(0.07)
    start_x = Inches(0.5)
    for i, (icon, titulo, desc, color, stat) in enumerate(nodes):
        x = start_x + (node_w + gap) * i
        add_rect(slide, x, top, node_w, node_h, fill=BG_CARD, border=color, border_w=2)
        add_textbox(slide, x, top + Inches(0.15), node_w, Inches(0.6),
                    icon, size=32, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(0.75), node_w, Inches(0.35),
                    titulo, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), top + Inches(1.1), node_w - Inches(0.2), Inches(0.65),
                    desc, size=9, color=TXT_BODY, align=PP_ALIGN.CENTER)

    # Stats reais embaixo
    add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2),
             fill=BG_CARD, border=EMERALD, border_w=1.5)
    add_textbox(slide, Inches(0.75), Inches(4.85), Inches(12), Inches(0.4),
                "📊  NÚMEROS REAIS DA VALIDAÇÃO DAS 5 SPRINTS",
                size=13, color=EMERALD_DARK, bold=True)

    stats = [
        ("322", "Testes",     "Playwright 100% OK"),
        ("79",  "UCs",        "17+13+19+11+19"),
        ("217", "RNs",        "153 enforçadas + 64 gap"),
        ("~480", "Screenshots", "Pares _acao + _resp"),
        ("5/5", "Sprints",    "Aprovadas com dados reais"),
        ("9/9", "DB checks",  "Verificação pós-teste"),
    ]
    x = Inches(0.75)
    for num, lbl, desc in stats:
        add_textbox(slide, x, Inches(5.3), Inches(2.0), Inches(0.55),
                    num, size=26, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.85), Inches(2.0), Inches(0.3),
                    lbl, size=11, color=TXT_INK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.15), Inches(2.0), Inches(0.6),
                    desc, size=9, color=TXT_MUTED, align=PP_ALIGN.CENTER)
        x += Inches(2.05)


def draw_timeline_sprints(slide):
    """Timeline horizontal das 5 sprints."""
    sprints_data = [
        (1, "Fundação",        "Empresa, portfolio, parametrizações"),
        (2, "Captação",        "PNCP, scores 6D, GO/NO-GO"),
        (3, "Precificação",    "6 camadas, proposta técnica"),
        (4, "Defesa",          "Impugnação, recursos, petições"),
        (5, "Pós-venda",       "Contratos, CRM 13 stages"),
    ]
    top = Inches(2.8)
    node_w = Inches(2.25)
    gap = Inches(0.3)
    start_x = Inches(0.75)
    # Linha base
    add_band(slide, Inches(0.9), top + Inches(1.25), Inches(11.5), Inches(0.04), EMERALD)
    for i, (num, nome, desc) in enumerate(sprints_data):
        x = start_x + (node_w + gap) * i
        color = SPRINT_COLORS[num]
        icon = SPRINT_ICONS[num]
        # Círculo do número
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), top + Inches(0.85),
                                        Inches(0.85), Inches(0.85))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.color.rgb = WHITE
        circ.line.width = Pt(3)
        tf = circ.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(num)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Card em cima ou embaixo alternado
        if i % 2 == 0:
            card_y = top - Inches(0.75)
        else:
            card_y = top + Inches(2.0)
        add_rect(slide, x, card_y, node_w, Inches(1.5), fill=BG_CARD, border=color, border_w=1.5)
        add_textbox(slide, x, card_y + Inches(0.1), node_w, Inches(0.4),
                    icon, size=22, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, card_y + Inches(0.55), node_w, Inches(0.35),
                    nome, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), card_y + Inches(0.9), node_w - Inches(0.2), Inches(0.55),
                    desc, size=9, color=TXT_MUTED, align=PP_ALIGN.CENTER)

    # Legenda embaixo
    add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.4),
                "✅  Todas as 5 sprints entregues e validadas  •  100% dos 322 testes Playwright aprovados",
                size=13, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)


def draw_leitura_ordem(slide):
    """Como o validador deve ler os documentos."""
    docs = [
        ("1", "Requisitos",           "requisitos_completosv8.md", "O QUE o sistema faz\n(RFs + 217 RNs)", ACCENT_BLUE),
        ("2", "Casos de uso V4",      "CASOS DE USO *.md",         "COMO o sistema usa\n(sequências + telas)", EMERALD),
        ("3", "Tutorial-2 (RP3X)",    "tutorialsprint#-2.md",       "GUIA passo-a-passo\n(o validador segue)", RGBColor(0x7C, 0x3A, 0xED)),
        ("4", "Resultado Validação",  "RESULTADO VALIDACAO #.md",   "PROVA de execução\n(screenshots + DB)", ACCENT_AMBER),
    ]
    top = Inches(2.4)
    card_w = Inches(2.95)
    card_h = Inches(3.9)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    for i, (num, titulo, arq, desc, color) in enumerate(docs):
        x = start_x + (card_w + gap) * i
        add_rect(slide, x, top, card_w, card_h, fill=BG_CARD, border=BORDER_SOFT)
        add_band(slide, x, top, card_w, Inches(0.35), color)
        add_textbox(slide, x, top + Inches(0.03), card_w, Inches(0.3),
                    num, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Título
        add_textbox(slide, x + Inches(0.15), top + Inches(0.5), card_w - Inches(0.3), Inches(0.45),
                    titulo, size=16, color=TXT_INK, bold=True)
        # Arquivo
        add_rect(slide, x + Inches(0.15), top + Inches(0.95), card_w - Inches(0.3), Inches(0.35),
                 fill=BG_CARD_ALT, border=BORDER_SOFT)
        add_textbox(slide, x + Inches(0.22), top + Inches(1.0), card_w - Inches(0.44), Inches(0.3),
                    arq, size=9, color=TXT_MUTED, font="Courier New")
        # Desc
        add_textbox(slide, x + Inches(0.15), top + Inches(1.5), card_w - Inches(0.3), Inches(1.2),
                    desc, size=11, color=TXT_BODY)
        # Ícone
        add_textbox(slide, x, top + Inches(2.8), card_w, Inches(0.9),
                    ["📋", "📘", "🗺️", "📊"][i], size=54, align=PP_ALIGN.CENTER)
        # Seta
        if i < len(docs) - 1:
            add_textbox(slide, x + card_w - Inches(0.1), top + Inches(1.7), gap + Inches(0.2), Inches(0.5),
                        "▶", size=18, color=TXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# EXCERPTS LITERAIS DOS DOCUMENTOS
# ═══════════════════════════════════════════════════════════

EXCERPTS = {
    "roadmap": {
        "title": "Roadmap fase 1 18-12-2025.docx",
        "source": "Roadmap fase 1 18-12-2025.docx · §1",
        "lines": [
            "Cadastro do Portfólio da empresa — Fase 1:",
            "Manuais técnicos, Instruções de Uso, Especificações",
            "Monitoramento das Fontes Públicas de Licitações",
        ],
    },
    "workflow": {
        "title": "WORKFLOW SISTEMA.pdf",
        "source": "WORKFLOW SISTEMA.pdf",
        "lines": [
            "Fluxo ponta-a-ponta: Descoberta → Análise → GO/NO-GO",
            "Precificação → Proposta → Submissão → Disputa",
            "Resultado → Contrato/Execução → CRM",
        ],
    },
    "sprint5_docx": {
        "title": "SPRINT 5 - VF.docx (fonte da sprint 5)",
        "source": "SPRINT 5 - VF.docx",
        "lines": [
            "Pós-venda: Execução de Contratos, Empenhos,",
            "Contratos a Vencer (5 tiers), CRM 13 stages,",
            "Parametrizações CRM, KPIs, Decisões",
        ],
    },
    "sprint_recursos_docx": {
        "title": "SPRINT RECURSOS E IMPUGNAÇÕES - V02.docx",
        "source": "SPRINT RECURSOS E IMPUGNAÇÕES - V02.docx",
        "lines": [
            "Validação legal do edital, esclarecimento/impugnação,",
            "Gerar petição, prazo, laudo, contra-razão,",
            "Submissão no portal — Lei 14.133/2021 art. 164-167",
        ],
    },
    "preco_proposta_3pdfs": {
        "title": "SPRINT PREÇO e PROPOSTA — 3 PDFs",
        "source": "PDF SPRINT PREÇO e PROPOSTA.pdf + ANOTAÇÕES REUNIÃO.pdf + REVISADA.pdf",
        "lines": [
            "1) PDF SPRINT PREÇO e PROPOSTA — escopo base Sprint 3",
            "2) ANOTAÇÕES REUNIÃO — lotes por especialidade, ERP/ICMS",
            "3) REVISADA — Edital → Seleção → Precificação → Proposta",
        ],
    },
    "requisitos_v8_rf": {
        "title": "requisitos_completosv8.md — RF-001",
        "source": "requisitos_completosv8.md, linha 63-76",
        "lines": [
            "RF-001: Cadastro da Empresa — Formulário com dados",
            "básicos (CNPJ, razão social, nome…), integração RFB,",
            "validação de DV, persistência na tabela empresa",
        ],
    },
    "requisitos_v8_rn": {
        "title": "requisitos_completosv8.md — RN-001",
        "source": "requisitos_completosv8.md, linha 2734",
        "lines": [
            "RN-001: O CNPJ da empresa é identificador único global",
            "— não pode existir mais de um cadastro com mesmo CNPJ.",
            "Enforçada em backend/models.py:Empresa.cnpj UNIQUE",
        ],
    },
    "requisitos_v8_camadas": {
        "title": "requisitos_completosv8.md — 3 camadas",
        "source": "requisitos_completosv8.md, linha 23-55",
        "lines": [
            "Fundação (Empresa, Portfolio, Parametrizações),",
            "Fluxo Comercial (11 etapas sequenciais),",
            "Painéis Transversais (CRM, KPIs, Alertas, Mapa)",
        ],
    },
    "uc_sprint1": {
        "title": "CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V4.md",
        "source": "CASOS DE USO ... PARAMETRIZACAO V4.md, linha 82-102",
        "lines": [
            "[UC-F01] Manter cadastro principal da empresa —",
            "Informações Cadastrais, Alertas IA, Histórico RF,",
            "edição inline via modal (RF-001, RF-005)",
        ],
    },
    "uc_sprint2": {
        "title": "CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V4.md",
        "source": "CASOS DE USO CAPTACAO ... V4.md, linha 135-142",
        "lines": [
            "[UC-CV01] Usuário clica em [Select] 'Selecione o Edital'",
            "Sistema verifica disponibilidade, retorna 12 resultados,",
            "ordena por score final (desc)",
        ],
    },
    "uc_sprint3": {
        "title": "CASOS DE USO PRECIFICACAO E PROPOSTA V4.md",
        "source": "CASOS DE USO PRECIFICACAO ... V4.md, linha 102-140",
        "lines": [
            "[UC-P01] Criar Lotes → Salvar Custos →",
            "Salvar Preço Base → Salvar Target →",
            "Salvar Lance Inicial → Salvar Lance Mínimo",
        ],
    },
    "uc_sprint4": {
        "title": "CASOS DE USO RECURSOS E IMPUGNACOES V4.md",
        "source": "CASOS DE USO RECURSOS ... V4.md, linha 92",
        "lines": [
            "✓ Tabela de inconsistências detectadas com Badges:",
            "ALTA (vermelho), MÉDIA (amarelo), BAIXA (verde)",
            "— e sugestão de ação (Esclarecimento/Impugnação)",
        ],
    },
    "uc_sprint5": {
        "title": "CASOS DE USO SPRINT5 V4.md",
        "source": "CASOS DE USO SPRINT5 V4.md, linha 31-40",
        "lines": [
            "[UC-FU01] Registrar Resultado, [UC-AT01] Buscar Atas,",
            "[UC-CT01] Cadastrar Contrato, [UC-CRM01] Pipeline 13 stages,",
            "[UC-CR01] Dashboard Contratado × Realizado",
        ],
    },
    "tutorial1_1": {
        "title": "tutorialsprint1-1.md (automação CH Hospitalar)",
        "source": "tutorialsprint1-1.md, linha 75-101",
        "lines": [
            "UC-F01: Navegar EmpresaPage → Preencher Razão Social,",
            "CNPJ, Endereço → Salvar → expect(tela).toContain('43.712.232')",
            "→ Screenshot F01_06_salvo.png (dados CH Hospitalar)",
        ],
    },
    "tutorial1_2": {
        "title": "tutorialsprint1-2.md (manual RP3X)",
        "source": "tutorialsprint1-2.md, linha 83-135",
        "lines": [
            "UC-F01 — O que este caso de uso faz: Aqui você",
            "vai preencher o cartão de identificação da empresa RP3X.",
            "✅ Correto se: Todos os campos preenchidos com os dados acima",
        ],
    },
    "tutorial2_2": {
        "title": "tutorialsprint2-2.md (manual RP3X)",
        "source": "tutorialsprint2-2.md, linha 114-123",
        "lines": [
            "Passo 2: Busca 1 — 'reagente hematologia',",
            "Score Rápido, Fonte PNCP.",
            "Você deve ver ~12 editais ordenados por score final",
        ],
    },
    "tutorial3_2": {
        "title": "tutorialsprint3-2.md (manual RP3X)",
        "source": "tutorialsprint3-2.md, linha 133-149",
        "lines": [
            "Lote 1: Hematologia (Itens 1-5),",
            "Lote 2: Bioquímica (Itens 6-7),",
            "Camada A Volumetria → calcular quantidade técnica",
        ],
    },
    "tutorial4_2": {
        "title": "tutorialsprint4-2.md (manual RP3X)",
        "source": "tutorialsprint4-2.md, linha 129-137",
        "lines": [
            "Inconsistências: cláusulas com ALTA, MÉDIA, BAIXA gravidade",
            "+ sugestão de ação (esclarecimento ou impugnação)",
            "Gera petição automática por IA (DeepSeek)",
        ],
    },
    "tutorial5_2": {
        "title": "tutorialsprint5-2.md (manual RP3X)",
        "source": "tutorialsprint5-2.md, linha 66-74",
        "lines": [
            "UC-CT07 — Gestão de Empenhos:",
            "Verificar EMPH-2026, Itens, Alerta SEM VALOR,",
            "Auditoria Empenho × Fatura × Entrega",
        ],
    },
    "resultado_sprint1": {
        "title": "RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        "source": "RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md, linha 8-52",
        "lines": [
            "✅ Resultado geral: 17/17 testes passaram (7.5 minutos)",
            "✓ Todos os campos cadastrais preenchidos e salvos",
            "Screenshots pareadas: F01_*_acao + F01_*_resp",
        ],
    },
    "resultado_sprint2": {
        "title": "RESULTADO VALIDACAO SPRINT2.md",
        "source": "RESULTADO VALIDACAO SPRINT2.md, linha 15-35",
        "lines": [
            "Total de UCs: 13, Testes aprovados: 18/18 (100%),",
            "Tempo: 16.7 minutos.",
            "UC-CV01 — Buscar editais por termo: 5 testes, APROVADO",
        ],
    },
    "resultado_sprint3": {
        "title": "RESULTADO VALIDACAO SPRINT3.md",
        "source": "RESULTADO VALIDACAO SPRINT3.md, linha 18-66",
        "lines": [
            "Total de UCs: 19, Testes aprovados: 19/19 (100%),",
            "Tempo: 9.7 minutos.",
            "✓ expect(element).toBeVisible() — 22 verificações",
        ],
    },
    "resultado_sprint4": {
        "title": "RESULTADO VALIDACAO SPRINT4.md",
        "source": "RESULTADO VALIDACAO SPRINT4.md, linha 16-46",
        "lines": [
            "Total de UCs: 14, Testes aprovados: 14/14 (100%),",
            "Screenshots: 57.",
            "✅ UC-FU01 — Registrar Resultado de Edital: APROVADO",
        ],
    },
    "resultado_sprint5": {
        "title": "RESULTADO VALIDACAO SPRINT5.md",
        "source": "RESULTADO VALIDACAO SPRINT5.md, linha 22-55",
        "lines": [
            "UCs V3 novos: 11, Testes aprovados: 11/11,",
            "Duração: 47.2s (paralelo, 10 workers).",
            "✓ assertDataVisible falha se tela não contém seletores",
        ],
    },
}


# ═══════════════════════════════════════════════════════════
# DIAGRAMAS NOVOS — V3
# ═══════════════════════════════════════════════════════════

def draw_pipeline_docs(slide):
    """Pipeline completo de documentos: fonte → requisitos → UC → tutoriais → execução → resultado."""
    # 6 estágios em flow horizontal
    stages = [
        ("📄", "FONTES",       "Roadmap, SPRINT docx,\nWORKFLOW SISTEMA",           ACCENT_BLUE),
        ("📘", "REQUISITOS",   "requisitos_completosv8\n60+ RFs · 217 RNs",          RGBColor(0x7C, 0x3A, 0xED)),
        ("📑", "CASOS DE USO", "CASOS DE USO V4 (5 docs)\n79 UCs especificados",    EMERALD),
        ("🗺️", "TUTORIAIS",    "tutorialsprint#-1 (CH)\ntutorialsprint#-2 (RP3X)",  ACCENT_AMBER),
        ("🤖", "EXECUÇÃO",     "Playwright 322 tests\n~480 screenshots",             RGBColor(0xDB, 0x27, 0x77)),
        ("✅", "RESULTADO",    "RESULTADO VALIDACAO\n5 relatórios · APROVADO",       EMERALD_DARK),
    ]
    top = Inches(2.35)
    card_w = Inches(2.05)
    card_h = Inches(2.85)
    gap = Inches(0.08)
    start_x = Inches(0.3)
    for i, (icon, titulo, desc, color) in enumerate(stages):
        x = start_x + (card_w + gap) * i
        add_rect(slide, x, top, card_w, card_h, fill=BG_CARD, border=color, border_w=2)
        add_band(slide, x, top, card_w, Inches(0.3), color)
        add_textbox(slide, x, top + Inches(0.02), card_w, Inches(0.28),
                    f"#{i+1}", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(0.38), card_w, Inches(0.6),
                    icon, size=36, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(1.05), card_w, Inches(0.4),
                    titulo, size=13, color=TXT_INK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.08), top + Inches(1.5), card_w - Inches(0.16), Inches(1.25),
                    desc, size=9, color=TXT_BODY, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            arrow_x = x + card_w - Inches(0.02)
            add_textbox(slide, arrow_x, top + Inches(1.1), gap + Inches(0.15), Inches(0.5),
                        "▶", size=14, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)

    # Bloco explicativo embaixo
    add_rect(slide, Inches(0.3), Inches(5.45), Inches(12.7), Inches(1.5),
             fill=EMERALD_LIGHT, border=EMERALD, border_w=1.5)
    add_textbox(slide, Inches(0.5), Inches(5.58), Inches(12.3), Inches(0.35),
                "💡  Cadeia de rastreabilidade — cada artefato se liga ao anterior",
                size=13, color=EMERALD_DARK, bold=True)
    add_textbox(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
                "Roadmap descreve a visão · SPRINT docx detalha escopo · requisitos_completosv8 consolida RFs/RNs · "
                "Casos de Uso V4 especificam cada tela · tutorial-1 vira roteiro de automação · tutorial-2 vira roteiro manual · "
                "Playwright executa contra CH Hospitalar · RESULTADO VALIDACAO prova aprovação · validador manual usa tutorial-2 + RP3X.",
                size=10, color=TXT_BODY)


def draw_jornada_docs(slide):
    """Jornada vertical: fonte → requisitos → UC → tutorial → teste → resultado em forma de escada."""
    layers = [
        ("1", "Fontes",       "Roadmap fase 1 18-12-2025.docx · WORKFLOW SISTEMA.pdf · SPRINT #.docx", ACCENT_BLUE),
        ("2", "Requisitos",   "requisitos_completosv8.md → 60+ RFs numerados + 217 RNs + 3 camadas",    RGBColor(0x7C, 0x3A, 0xED)),
        ("3", "Casos de uso", "5 arquivos CASOS DE USO ... V4.md → 79 UCs (atores, pré, sequência…)",    EMERALD),
        ("4", "Tutoriais",    "tutorialsprint#-1 (automação CH) · tutorialsprint#-2 (manual RP3X)",      ACCENT_AMBER),
        ("5", "Execução",     "Playwright tests (322) · seeds · helpers · screenshots pareadas",        RGBColor(0xDB, 0x27, 0x77)),
        ("6", "Resultado",    "5× RESULTADO VALIDACAO SPRINT#.md · PPT comercial · PDFs · PR",          EMERALD_DARK),
    ]
    top = Inches(2.25)
    row_h = Inches(0.78)
    gap = Inches(0.06)
    for i, (num, titulo, desc, color) in enumerate(layers):
        y = top + (row_h + gap) * i
        # Card horizontal
        add_rect(slide, Inches(0.6), y, Inches(12.1), row_h, fill=BG_CARD, border=BORDER_SOFT)
        # Faixa lateral
        add_band(slide, Inches(0.6), y, Inches(0.18), row_h, color)
        # Número em círculo
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.15),
                                       Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.color.rgb = WHITE; circ.line.width = Pt(2)
        tf = circ.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = num; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE
        # Título
        add_textbox(slide, Inches(1.65), y + Inches(0.12), Inches(3.0), Inches(0.3),
                    titulo, size=15, color=color, bold=True)
        # Desc
        add_textbox(slide, Inches(1.65), y + Inches(0.42), Inches(11.0), Inches(0.35),
                    desc, size=11, color=TXT_BODY)
        # Seta entre linhas
        if i < len(layers) - 1:
            add_textbox(slide, Inches(0.92), y + row_h - Inches(0.05), Inches(0.6), gap + Inches(0.15),
                        "▼", size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def draw_pipeline_ia_dados(slide):
    """Pipeline de dados: PNCP + DeepSeek + Brave + ANVISA → banco → UI → validador."""
    nodes = [
        ("🌐", "PNCP",     "Portal Nacional\nde Contratações",      ACCENT_BLUE),
        ("🤖", "DeepSeek", "LLM via\nOpenRouter",                    RGBColor(0x7C, 0x3A, 0xED)),
        ("🔎", "Brave",    "Busca web\nde contexto",                 ACCENT_AMBER),
        ("🩺", "ANVISA",   "Registro de\nprodutos",                  EMERALD),
    ]
    # 4 fontes à esquerda
    top = Inches(2.3)
    for i, (icon, nome, desc, color) in enumerate(nodes):
        y = top + (Inches(1.05)) * i
        add_rect(slide, Inches(0.5), y, Inches(3.0), Inches(0.9),
                 fill=BG_CARD, border=color, border_w=1.5)
        add_textbox(slide, Inches(0.6), y + Inches(0.15), Inches(0.6), Inches(0.6),
                    icon, size=24, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(2.0), Inches(0.35),
                    nome, size=13, color=color, bold=True)
        add_textbox(slide, Inches(1.3), y + Inches(0.42), Inches(2.0), Inches(0.45),
                    desc, size=9, color=TXT_BODY)

    # Backend central
    add_rect(slide, Inches(4.5), Inches(2.5), Inches(3.5), Inches(3.6),
             fill=EMERALD_LIGHT, border=EMERALD, border_w=2)
    add_textbox(slide, Inches(4.5), Inches(2.65), Inches(3.5), Inches(0.5),
                "⚙️  Backend", size=18, color=EMERALD_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.5), Inches(3.1), Inches(3.5), Inches(0.35),
                "Flask · port 5007", size=11, color=EMERALD_DARK, align=PP_ALIGN.CENTER)
    backend_items = [
        "app.py (main API)",
        "crud_routes.py",
        "empenho_routes.py",
        "crm_routes.py",
        "tools.py (DeepSeek)",
        "scheduler.py (jobs)",
    ]
    y = Inches(3.45)
    for it in backend_items:
        add_textbox(slide, Inches(4.7), y, Inches(3.3), Inches(0.3),
                    "• " + it, size=10, color=TXT_BODY)
        y += Inches(0.3)
    add_textbox(slide, Inches(4.5), Inches(5.7), Inches(3.5), Inches(0.35),
                "SQLAlchemy + MySQL", size=10, color=EMERALD_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Frontend à direita
    add_rect(slide, Inches(9.0), Inches(2.5), Inches(3.9), Inches(3.6),
             fill=BG_CARD, border=ACCENT_BLUE, border_w=2)
    add_textbox(slide, Inches(9.0), Inches(2.65), Inches(3.9), Inches(0.5),
                "🖥️  Frontend", size=18, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.0), Inches(3.1), Inches(3.9), Inches(0.35),
                "React · Vite · port 5180", size=11, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    frontend_items = [
        "CaptacaoPage · ValidacaoPage",
        "PrecificacaoPage · PropostaPage",
        "RecursosPage · FollowupPage",
        "ProducaoPage · CRMPage",
        "KPIsPage · MapaPage",
    ]
    y = Inches(3.55)
    for it in frontend_items:
        add_textbox(slide, Inches(9.15), y, Inches(3.7), Inches(0.3),
                    "• " + it, size=10, color=TXT_BODY)
        y += Inches(0.32)

    # Setas
    add_textbox(slide, Inches(3.55), Inches(4.0), Inches(0.9), Inches(0.5),
                "▶", size=24, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.05), Inches(4.0), Inches(0.9), Inches(0.5),
                "▶", size=24, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)

    # Validador embaixo
    add_rect(slide, Inches(0.5), Inches(6.25), Inches(12.4), Inches(0.75),
             fill=BG_CARD, border=EMERALD_DARK, border_w=1.5)
    add_textbox(slide, Inches(0.7), Inches(6.38), Inches(0.6), Inches(0.5),
                "👤", size=22, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), Inches(6.32), Inches(11.4), Inches(0.3),
                "Validador humano — valida2@valida.com.br",
                size=12, color=EMERALD_DARK, bold=True)
    add_textbox(slide, Inches(1.4), Inches(6.6), Inches(11.4), Inches(0.3),
                "Usa tutorialsprint#-2.md com a empresa RP3X e reporta divergências via ChatGPT→Claude",
                size=10, color=TXT_BODY)


def draw_excerpt_card(slide, excerpt_key, left, top, width, height,
                      *, color=EMERALD, bold_title=True):
    """Renderiza um card de excerpt literal a partir da chave EXCERPTS[key]."""
    exc = EXCERPTS.get(excerpt_key)
    if not exc:
        return
    add_rect(slide, left, top, width, height,
             fill=EMERALD_LIGHT, border=color, border_w=1.5)
    # Título
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.4), Inches(0.35),
                f"📄 {exc['title']}", size=12, color=EMERALD_DARK, bold=bold_title)
    # Linhas literais (monospace pra dar cara de código)
    y = top + Inches(0.5)
    body_h = height - Inches(0.95)
    line_h = Inches(0.28)
    for line in exc["lines"]:
        add_textbox(slide, left + Inches(0.3), y,
                    width - Inches(0.5), line_h,
                    "› " + line, size=10, color=TXT_BODY, font="Courier New")
        y += line_h
    # Fonte
    add_textbox(slide, left + Inches(0.2), top + height - Inches(0.32),
                width - Inches(0.4), Inches(0.25),
                f"Fonte: {exc['source']}", size=8, color=TXT_MUTED)


def excerpts_showcase_slide(prs, *, title, kicker, subtitle, excerpt_keys,
                            footer=None, sprint_num=None):
    """Slide que mostra 2-4 excerpts lado-a-lado em grid."""
    slide = new_slide(prs)
    slide_header(slide, kicker, title, subtitle, sprint_num=sprint_num)
    n = len(excerpt_keys)
    if n <= 2:
        card_w = Inches(6.0); card_h = Inches(4.8)
        cols = n; rows = 1
    elif n <= 4:
        card_w = Inches(6.0); card_h = Inches(2.35)
        cols = 2; rows = 2
    elif n <= 6:
        card_w = Inches(4.0); card_h = Inches(2.35)
        cols = 3; rows = 2
    else:
        card_w = Inches(4.0); card_h = Inches(1.75)
        cols = 3; rows = (n + 2) // 3
    gap = Inches(0.2)
    total_w = card_w * cols + gap * (cols - 1)
    start_x = (Inches(13.333) - total_w) / 2
    start_y = Inches(2.1)
    for i, key in enumerate(excerpt_keys):
        col = i % cols; row = i // cols
        x = start_x + (card_w + gap) * col
        y = start_y + (card_h + gap) * row
        color = SPRINT_COLORS[sprint_num] if sprint_num else EMERALD
        draw_excerpt_card(slide, key, x, y, card_w, card_h, color=color)
    add_footer(slide, footer or "Trechos literais extraídos dos documentos originais",
               sprint_color=SPRINT_COLORS.get(sprint_num))
    return slide


def content_slide_with_excerpt(prs, *, kicker, title, subtitle, bullets,
                                screenshot_uc=None, excerpt_key=None,
                                callouts_rel=None, footer=None, sprint_num=None):
    """Variação do content_slide que inclui um card de excerpt no canto.

    callouts_rel: lista de (x_rel, y_rel, w_rel, h_rel, explicacao) — callouts
                  numerados em cima da screenshot (coordenadas 0-1).
    """
    slide = new_slide(prs)
    slide_header(slide, kicker, title, subtitle, sprint_num=sprint_num)

    has_screenshot = bool(screenshot_uc)
    img_path = screenshot_path(screenshot_uc) if screenshot_uc else None

    # Layout idêntico ao content_slide normal
    bullets_left  = Inches(0.6)
    bullets_width = Inches(5.2)
    shot_left     = Inches(6.0)
    shot_top      = Inches(2.1)
    shot_width    = Inches(6.8)
    shot_height   = Inches(3.55)

    add_rect(slide, bullets_left, Inches(2.05), bullets_width, Inches(3.65),
             fill=BG_CARD, border=BORDER_SOFT)
    y = Inches(2.25)
    color = SPRINT_COLORS[sprint_num] if sprint_num else EMERALD
    for b in bullets:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bullets_left + Inches(0.25),
                                       y + Inches(0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, bullets_left + Inches(0.5), y,
                    bullets_width - Inches(0.6), Inches(0.7),
                    b, size=11, color=TXT_BODY)
        y += Inches(0.6)

    if has_screenshot:
        add_screenshot(slide, img_path, shot_left, shot_top, shot_width, shot_height)
        if callouts_rel and img_path:
            for i, co in enumerate(callouts_rel, start=1):
                x_rel, y_rel, w_rel, h_rel, _ = co
                hl_left = shot_left + Emu(int(shot_width * x_rel))
                hl_top  = shot_top + Emu(int(shot_height * y_rel))
                hl_w    = Emu(int(shot_width * w_rel))
                hl_h    = Emu(int(shot_height * h_rel))
                add_highlight_box(slide, hl_left, hl_top, hl_w, hl_h)
                add_circle_num(slide, hl_left - Inches(0.15), hl_top - Inches(0.15),
                               Inches(0.35), i)
            # Legenda numerada abaixo do screenshot
            parts = [f"({i}) {co[4]}" for i, co in enumerate(callouts_rel, start=1)]
            add_textbox(slide, shot_left, shot_top + shot_height + Inches(0.1),
                        shot_width, Inches(0.3),
                        "   ".join(parts), size=9, color=TXT_MUTED)

    # Excerpt no fundo largura total
    if excerpt_key:
        draw_excerpt_card(slide, excerpt_key,
                          Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.15),
                          color=color)

    add_footer(slide, footer or "", sprint_color=SPRINT_COLORS.get(sprint_num))
    return slide


# ═══════════════════════════════════════════════════════════
# BUILD DA APRESENTAÇÃO
# ═══════════════════════════════════════════════════════════


# ═══════════════════════════════════════════════════════════
# ADAPTADORES: API make_* → API V2
# ═══════════════════════════════════════════════════════════
# Mantém 83-slide structure do gerador recuperado, mas roteia
# cada make_* para cover_slide/content_slide/sprint_cover do V2.

_SECTION_KICKER = [None]  # mutable holder

def _set_section_kicker(text):
    _SECTION_KICKER[0] = text

def make_cover_slide(prs):
    return cover_slide(prs)

def make_sprint_cover(prs, sprint_num, title, subtitle, bullets):
    _set_section_kicker(None)  # reseta kicker ao trocar de sprint
    return sprint_cover(prs, sprint_num, title, subtitle, bullets)

def _kicker_for(sprint_num, subtitle):
    """Gera kicker curto. Usa última seção setada, ou subtítulo truncado, ou sprint."""
    if _SECTION_KICKER[0]:
        return _SECTION_KICKER[0]
    if sprint_num:
        names = {1: "SPRINT 1 · FUNDAÇÃO", 2: "SPRINT 2 · CAPTAÇÃO",
                 3: "SPRINT 3 · PRECIFICAÇÃO", 4: "SPRINT 4 · DEFESA",
                 5: "SPRINT 5 · PÓS-VENDA"}
        return names.get(sprint_num, f"SPRINT {sprint_num}")
    # fallback: primeira palavra/frase do subtitle maiúscula
    if subtitle:
        s = subtitle.split('|')[0].split('—')[0].strip()
        s = s.split(':')[0].strip()
        return s.upper()[:40] if s else "VISÃO GERAL"
    return "VISÃO GERAL"

def make_content_slide(prs, title, subtitle, bullets, screenshot_uc=None,
                       footer=None, sprint_num=None):
    kicker = _kicker_for(sprint_num, subtitle)
    return content_slide(
        prs,
        kicker=kicker,
        title=title,
        subtitle=subtitle,
        bullets=bullets,
        screenshot_uc=screenshot_uc,
        footer=footer,
        sprint_num=sprint_num,
    )



def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ═══════════════════════════════════════════════════
    # SEÇÃO A: Capa e Visão Geral (slides 1-5)
    # ═══════════════════════════════════════════════════

    # Slide 1 — Capa
    make_cover_slide(prs)

    # Slide 2 — A jornada
    make_content_slide(prs,
        "A jornada de uma licitação",
        "O sistema cobre TODAS as etapas do ciclo de vida",
        [
            "A empresa existe no sistema (Sprint 1: cadastro, produtos, configurações)",
            "Ela descobre um edital interessante (Sprint 2: captação e validação com IA)",
            "Ela calcula o preço e monta a proposta (Sprint 3: 6 camadas de precificação)",
            "Ela se defende juridicamente se necessário (Sprint 4: impugnação, recursos)",
            "Ela acompanha resultados e gerencia contratos (Sprint 5: follow-up, CRM)",
        ])

    # Slide 3 — Organização
    make_content_slide(prs,
        "Como o sistema está organizado",
        "3 camadas que se complementam",
        [
            "FUNDAÇÃO: empresa, produtos, certificações, configurações (Sprint 1)",
            "FLUXO COMERCIAL: 11 etapas sequenciais do edital ao contrato",
            "PAINÉIS TRANSVERSAIS: CRM, KPIs, alertas, relatórios",
            "Cada sprint construiu uma parte deste quebra-cabeça",
            "As sprints são sequenciais — Sprint 2 depende da Sprint 1, e assim por diante",
        ])

    # Slide 4 — Em números
    make_content_slide(prs,
        "Em números",
        "O que foi entregue nas 5 sprints",
        [
            "79 casos de uso especificados e implementados",
            "60+ requisitos funcionais rastreados (RF-001 a RF-060)",
            "5 sprints entregues com 100% de aprovação nos testes automatizados",
            "322 testes Playwright executados com sucesso",
            "2 empresas de teste: CH Hospitalar (automação) + RP3X (validação manual)",
        ])

    # Slide 5 — Empresas
    make_content_slide(prs,
        "As duas empresas de teste",
        "Cada empresa tem um papel diferente na validação",
        [
            "CH Hospitalar — equipamentos médicos hospitalares (CNPJ 43.712.232)",
            "  → Usada nos testes automatizados (Playwright) — tutorial-1",
            "RP3X — reagentes para diagnóstico in vitro (CNPJ 68.218.593)",
            "  → Usada na validação manual pelo validador — tutorial-2",
            "Credenciais: valida1@valida.com.br (CH) / valida2@valida.com.br (RP3X) — senha: 123456",
        ])

    # ═══════════════════════════════════════════════════
    # SEÇÃO A.5 [V3]: Pipeline de Documentos
    # ═══════════════════════════════════════════════════

    _set_section_kicker("DOCUMENTOS · PIPELINE")

    # Slide novo — Pipeline completo de documentos (diagrama horizontal)
    diagram_slide(prs,
        kicker="DOCUMENTOS · PIPELINE",
        title="O pipeline de documentos",
        subtitle="Da visão inicial ao relatório de aprovação — cadeia completa de rastreabilidade",
        draw_fn=draw_pipeline_docs,
        footer="Fontes originais → Requisitos → Casos de Uso → Tutoriais → Execução → Resultado")

    # Slide novo — Jornada vertical dos docs
    diagram_slide(prs,
        kicker="DOCUMENTOS · JORNADA",
        title="Jornada vertical dos documentos",
        subtitle="Mesma cadeia vista como escada — cada camada depende da anterior",
        draw_fn=draw_jornada_docs,
        footer="Cada doc se liga ao próximo por rastreabilidade — RF→UC→Teste→Resultado")

    # Slide novo — Fontes originais (Roadmap, WORKFLOW, SPRINT docx, PDFs Preço/Proposta)
    excerpts_showcase_slide(prs,
        kicker="DOCUMENTOS · FONTES",
        title="As fontes originais do sistema",
        subtitle="Roadmap, WORKFLOW, SPRINT docx e PDFs de Preço/Proposta — onde tudo começou",
        excerpt_keys=["roadmap", "workflow", "sprint5_docx",
                      "sprint_recursos_docx", "preco_proposta_3pdfs"],
        footer="Documentos entregues pelo dono do projeto antes da Sprint 1")

    _set_section_kicker(None)

    # ═══════════════════════════════════════════════════
    # SEÇÃO B: Processo de Validação (slides 6-12)
    # ═══════════════════════════════════════════════════

    make_content_slide(prs,
        "O ciclo de validação: 3 documentos por sprint",
        "Para cada sprint, existem 3 tipos de documento",
        [
            "CASOS DE USO — O que o sistema DEVE fazer (especificação detalhada)",
            "TUTORIAL-2 — Roteiro para o VALIDADOR testar manualmente (dados RP3X)",
            "RESULTADO DA VALIDAÇÃO — O que a automação Playwright comprovou (screenshots CH Hospitalar)",
            "Além disso: requisitos_completosv7.md (documento mestre de 60+ RFs)",
            "O validador recebe tudo pronto: especificação + roteiro + evidências",
        ])

    make_content_slide(prs,
        "O documento de Requisitos (requisitos_completosv8.md)",
        "O documento mestre — cada funcionalidade tem um código RF-XXX",
        [
            "Cada funcionalidade do sistema tem um código: RF-001, RF-019, RF-039...",
            "Exemplo: RF-019 = Buscar editais no PNCP com scores automáticos",
            "São 60+ requisitos organizados em 3 camadas (Fundação, Fluxo, Painéis)",
            "Cada caso de uso implementa um ou mais RFs — rastreabilidade completa",
            "217 Regras de Negócio (RNs) anotadas — 153 enforçadas + 64 em plano",
            "Também disponível em PDF: docs/requisitos_completosv8.pdf",
        ])

    # Slide novo V3 — Excerpts literais do requisitos_completosv8
    excerpts_showcase_slide(prs,
        kicker="REQUISITOS · v8",
        title="Trechos literais — requisitos_completosv8",
        subtitle="Exemplos reais de como um RF, uma RN e as 3 camadas aparecem no documento",
        excerpt_keys=["requisitos_v8_rf", "requisitos_v8_rn", "requisitos_v8_camadas"],
        footer="requisitos_completosv8.md (3.204 linhas) · versão PDF disponível")

    make_content_slide(prs,
        "O que é um Caso de Uso?",
        "Descreve UMA funcionalidade específica do sistema",
        [
            "É a UNIDADE de validação — você valida UM UC por vez, não a sprint inteira",
            "ATOR: quem usa (analista comercial, gestor, fiscal de contrato...)",
            "PRÉ-CONDIÇÕES: o que precisa existir antes (ex: empresa cadastrada com CNPJ)",
            "SEQUÊNCIA DE EVENTOS: passos numerados — cada clique e cada resposta do sistema",
            "TELA REPRESENTATIVA: layout hierárquico mostrando cada botão, campo e badge",
            "ASSERTIONS: o que TEM que aparecer pra você dizer 'passou' (ex: tabela ≥3 linhas)",
            "REGRAS DE NEGÓCIO (RNs): as leis internas que o UC deve respeitar (ex: RN-049)",
            "Cada UC tem código único: UC-F01 (Sprint 1), UC-CV08 (Sprint 2), UC-CRM01 (Sprint 5)",
        ])

    make_content_slide(prs,
        "Onde encontrar os Casos de Uso",
        "5 documentos V4 — um por sprint — todos na pasta docs/",
        [
            "Sprint 1 → CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V4.md (17 UCs: F01..F17)",
            "Sprint 2 → CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V4.md (13 UCs: CV01..CV13)",
            "Sprint 3 → CASOS DE USO PRECIFICACAO E PROPOSTA V4.md (19 UCs: P01..P12 + R01..R07)",
            "Sprint 4 → CASOS DE USO RECURSOS E IMPUGNACOES V4.md (11 UCs: I01..I05 + RE01..RE06)",
            "Sprint 5 → CASOS DE USO SPRINT5 V4.md (19 UCs: FU + AT + CT + CRM + CR)",
            "A versão V4 é a mais recente — traz cada UC com as RNs aplicáveis listadas",
            "Todos os 217 RNs (RN-001..RN-217) estão anotados nos UCs pertinentes",
        ])

    # Slide novo V3 — Trechos literais dos 5 Casos de Uso V4
    excerpts_showcase_slide(prs,
        kicker="CASOS DE USO · V4",
        title="Trechos literais dos Casos de Uso V4",
        subtitle="Um exemplo de cada sprint: como o documento descreve um UC",
        excerpt_keys=["uc_sprint1", "uc_sprint2", "uc_sprint3", "uc_sprint4", "uc_sprint5"],
        footer="5 arquivos CASOS DE USO ... V4.md · 79 UCs totais")

    make_content_slide(prs,
        "O Tutorial-2: seu roteiro de validação manual",
        "Cada sprint tem um arquivo tutorial-sprint[N]-2.md com o passo-a-passo da RP3X",
        [
            "É um SCRIPT pronto: 'Clique no botão X, digite Y, espere a tela mostrar Z'",
            "Empresa de teste: RP3X (valida2@valida.com.br / 123456) — diferente da CH",
            "TODOS os dados que você precisa digitar JÁ ESTÃO no tutorial — não invente",
            "Cada passo tem uma 'evidência esperada' — se diferir, é um achado",
            "Execute na ORDEM: Sprint 1 → 2 → 3 → 4 → 5 (uma sprint depende da anterior)",
            "Tutorial-1 (CH Hospitalar) já foi executado por Playwright — ele é a referência;",
            "  o tutorial-2 (RP3X) é o que VOCÊ executa manualmente",
        ])

    # Slide novo V3 — Tutoriais lado-a-lado (1 e 2)
    excerpts_showcase_slide(prs,
        kicker="TUTORIAIS · 1 + 2",
        title="Tutoriais — trechos literais",
        subtitle="Tutorial-1 vira script Playwright (CH Hospitalar) · Tutorial-2 é o roteiro manual (RP3X)",
        excerpt_keys=["tutorial1_1", "tutorial1_2", "tutorial2_2",
                      "tutorial3_2", "tutorial4_2", "tutorial5_2"],
        footer="10 arquivos de tutorial · 5 sprints × 2 versões (automação + manual)")

    make_content_slide(prs,
        "A validação automatizada (Playwright)",
        "Testes automáticos JÁ EXECUTADOS com a empresa CH Hospitalar",
        [
            "O sistema Playwright executou os tutorial-1 automaticamente",
            "Para cada passo: 2 screenshots — AÇÃO (o que foi clicado) e RESPOSTA (resultado)",
            "Se a tela estivesse vazia, o teste REPROVARIA automaticamente",
            "Helper assertDataVisible: garante que há dados reais antes do screenshot",
            "Todos os 322 testes passaram com 100% de aprovação",
        ])

    make_content_slide(prs,
        "O Resultado da Validação",
        "Evidência de que cada UC funciona corretamente",
        [
            "Para cada sprint existe um documento: RESULTADO VALIDACAO SPRINT[N].md",
            "Tabela de cada UC: código, título, passos, screenshots, status (APROVADO)",
            "Screenshots da automação são a PROVA de que o sistema funciona",
            "Neste PowerPoint, cada slide de UC mostra a screenshot da validação real",
            "O validador pode comparar o que vê no tutorial-2 com estas evidências",
        ])

    make_content_slide(prs,
        "O Relatório de Aceitação",
        "Parecer final de cada sprint",
        [
            "Lista todos os requisitos funcionais (RF) atendidos",
            "Evidências: código dos testes, screenshots, seed de dados",
            "Conclusão: APROVADO INTEGRALMENTE em todas as 5 sprints",
            "Pode ser usado como checklist pelo validador",
            "Arquivos: relatorio_aceitacao_sprint[N].md",
        ])

    # Slide novo V3 — Trechos literais dos RESULTADO VALIDACAO
    excerpts_showcase_slide(prs,
        kicker="RESULTADOS · APROVADO",
        title="Trechos dos relatórios de validação — as 5 sprints",
        subtitle="Cada linha é cópia literal do respectivo RESULTADO VALIDACAO SPRINT#.md",
        excerpt_keys=["resultado_sprint1", "resultado_sprint2",
                      "resultado_sprint3", "resultado_sprint4", "resultado_sprint5"],
        footer="5 relatórios · total: 322 testes Playwright · 100% aprovados")

    # ═══════════════════════════════════════════════════
    # SEÇÃO B.1: Exemplo real de Caso de Uso (UC-CV08)
    # ═══════════════════════════════════════════════════

    make_content_slide(prs,
        "Anatomia de um Caso de Uso — Exemplo real",
        "UC-CV08 (Sprint 2): Calcular scores multidimensionais e decidir GO/NO-GO",
        [
            "ATOR: Analista comercial",
            "RFs RELACIONADOS: RF-027 (scores 6 dimensões), RF-028 (decisão IA), RF-037 (justificativa)",
            "RNs APLICADAS (15): RN-047 (peso técnico), RN-048 (peso comercial), RN-049 (limiar GO),",
            "  RN-050 (limiar NO-GO), RN-052 (cache scores), RN-080 (versionar decisão)... etc.",
            "PRÉ-CONDIÇÕES: edital salvo, itens importados, pesos GO/NO-GO configurados (UC-F14)",
            "PÓS-CONDIÇÕES: scores persistidos, decisão registrada com timestamp e usuário,",
            "  edital muda de status para 'em_analise' (RN-082)",
        ],
        footer="docs/CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V4.md — UC-CV08")

    make_content_slide(prs,
        "UC-CV08 — Sequência de eventos (trecho literal do documento)",
        "Cada passo numerado: ação do ator + resposta do sistema",
        [
            "1. Usuário abre [Aba: 'Aderência'] no Painel de Abas da ValidacaoPage",
            "2. Usuário clica em [Botão: 'Calcular Scores IA']",
            "3. Sistema chama POST /api/editais/{id}/scores-validacao (DeepSeek)",
            "4. Sistema retorna 6 scores (0-100) + score_final + decisao_sugerida + justificativa",
            "5. Tela atualiza: ScoreCircle central + 6 ScoreBars + Badge 'Decisão IA: GO'",
            "6. Usuário lê justificativa em [Box: 'Pontos Positivos / Pontos de Atenção']",
            "7. Usuário escolhe em [Radio]: GO / Acompanhar / NO-GO (sobrepõe sugestão IA)",
            "8. Usuário escreve justificativa própria em [TextArea] e clica [Botão: 'Salvar']",
            "→ Sistema persiste em DecisaoGoNoGoHistorico + atualiza Edital.status",
        ],
        footer="RN-049: GO exige score_final≥70 + tec≥60 + jur≥60 | RN-080: toda decisão vira histórico")

    make_content_slide(prs,
        "UC-CV08 — Assertions (como o validador sabe que passou)",
        "Lista de evidências obrigatórias na tela após a execução",
        [
            "✓ ScoreCircle central com valor numérico (ex: 78/100) e cor (verde/amarelo/vermelho)",
            "✓ 6 ScoreBars preenchidas, cada uma com nome e valor:",
            "    Técnica, Documental, Complexidade, Jurídico, Logístico, Comercial",
            "✓ Badge 'Decisão IA' visível: GO (verde) / Acompanhar (amarelo) / NO-GO (vermelho)",
            "✓ Pontos Positivos (verde) e Pontos de Atenção (amarelo) com texto NÃO-vazio",
            "✓ Após salvar: toast 'Decisão registrada' e linha nova no histórico",
            "✗ Se algo estiver vazio, em inglês, ou genérico → REPROVAR e reportar ao Claude",
        ])

    # ═══════════════════════════════════════════════════
    # SEÇÃO B.2: Ordem de leitura dos documentos
    # ═══════════════════════════════════════════════════

    make_content_slide(prs,
        "Ordem de leitura — seu fluxo de validação por sprint",
        "Para CADA sprint, sempre nesta sequência: teoria → evidência → mão-na-massa",
        [
            "1º) CASOS DE USO V4 da sprint  — leia TODOS os UCs da sprint",
            "       'O que o sistema DEVE fazer? Quais RNs aplicam? Quais assertions?'",
            "2º) RESULTADO VALIDACAO SPRINT[N].md  — veja o que a automação já provou",
            "       'Como ficou a tela quando a CH Hospitalar passou pelos mesmos passos?'",
            "3º) TUTORIAL-SPRINT[N]-2.md  — siga o roteiro com a empresa RP3X",
            "       'Agora EU faço cada passo e comparo com o que vi no resultado da CH'",
            "Por quê nesta ordem? UC ensina a teoria, resultado mostra a prática,",
            "  tutorial é mão-na-massa. Pular etapas = você valida sem entender o porquê.",
        ])

    make_content_slide(prs,
        "Exemplo concreto da ordem de leitura — Sprint 2",
        "Como o validador prepara a Sprint 2 antes de abrir o sistema",
        [
            "PASSO 1 (30 min): abra docs/CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V4.md",
            "  Leia UC-CV01 (buscar no PNCP) → UC-CV08 (scores GO/NO-GO) → UC-CV13 (chat IA)",
            "  Anote: 'UC-CV08 espera 6 ScoreBars + decisão da IA + justificativa'",
            "PASSO 2 (15 min): abra docs/RESULTADO VALIDACAO SPRINT2.md",
            "  Veja a screenshot CV08_scores.png da CH Hospitalar — confirme que tem 6 barras",
            "PASSO 3 (executar): abra docs/tutorial-sprint2-2.md, faça login com RP3X",
            "  Siga passo a passo: 'Buscar hemograma → salvar edital → calcular scores → decidir GO'",
            "  Compare a SUA tela com a screenshot da CH — divergência = achado a reportar",
        ])

    # ═══════════════════════════════════════════════════
    # SEÇÃO B.3: Como registrar os resultados (ChatGPT → Claude)
    # ═══════════════════════════════════════════════════

    make_content_slide(prs,
        "O CICLO de registro de observações — visão geral",
        "Você nunca escreve direto pro Claude. Sempre passa pelo ChatGPT primeiro",
        [
            "Por quê? Porque o ChatGPT transforma sua FALA INFORMAL em PROMPT ESTRUTURADO",
            "  que o Claude Code entende com mínima ambiguidade.",
            "FLUXO COMPLETO em 4 passos:",
            "  PASSO 1 — CAPTURAR: você fala/digita observações soltas no ChatGPT",
            "  PASSO 2 — ESTRUTURAR: pede ao ChatGPT pra virar prompt formal pro Claude",
            "  PASSO 3 — COLAR: copia a saída do ChatGPT pra um arquivo .md",
            "  PASSO 4 — ENTREGAR: envia o .md ao Claude Code via chat do Anthropic",
            "Resultado: o Claude recebe instruções claras, sem perder contexto, sem ruído",
        ])

    make_content_slide(prs,
        "Passo 1 — CAPTURAR observações no ChatGPT (durante o teste)",
        "Mantenha uma janela do ChatGPT aberta enquanto executa o tutorial-2",
        [
            "Abra o ChatGPT em uma aba SEPARADA do navegador onde você está testando",
            "Comece o chat com: 'Estou validando a Sprint 2 do Facilitia.ia. Vou narrar",
            "  achados pra você ir guardando. No final eu peço pra você estruturar.'",
            "A CADA tela que você testa, narre o que viu — pode ser informal:",
            "  'Tela do UC-CV08, calculei o score, deu 78 mas a Decisão IA tá em branco.",
            "   Os pontos positivos vieram em inglês. Tirei screenshot e salvei no desktop.'",
            "Se quiser, use a função de DITAR do ChatGPT — fala em vez de digitar",
            "Não pare pra formatar — narre tudo no fluxo, é mais rápido",
        ])

    make_content_slide(prs,
        "Passo 1 — Exemplo REAL de captura informal",
        "O que você digita no ChatGPT enquanto testa a Sprint 2",
        [
            "[texto que você cola no ChatGPT, tal qual fala]",
            "",
            "'Sprint 2, RP3X. Cliquei em Captação, busquei hemograma, voltou 12 editais",
            "ok. Salvei o primeiro. Aí abri Validação, escolhi o edital, fui na aba",
            "Aderência. Cliquei Calcular Scores. Demorou uns 15 segundos, score deu 78",
            "mas a barra Jurídico ficou só 30. Mesmo assim a IA recomendou GO. Estranho.",
            "",
            "Aí no UC-CV13, o chat. Perguntei resuma este edital. Veio resposta boa em PT.",
            "Mas perguntei quais riscos e veio só que não tem riscos, sem detalhar nada.",
            "",
            "Por fim no UC-CV12 análise de mercado, o gráfico de pizza não carregou,",
            "ficou só um quadrado branco. Tirei print, tá em prints/sprint2/cv12.png'",
        ])

    make_content_slide(prs,
        "Passo 2 — ESTRUTURAR: peça ao ChatGPT pra virar prompt do Claude",
        "Quando terminar a sprint, peça a transformação com este texto exato",
        [
            "[copie e cole no ChatGPT, depois de ter narrado tudo]",
            "",
            "'Agora transforme TUDO que eu narrei acima em um PROMPT BEM ESTRUTURADO",
            "para o Claude Code (Anthropic). Use este formato para cada achado:",
            "",
            "  [UC-XXX] SEVERIDADE (CRÍTICA/ALTA/MÉDIA/BAIXA)",
            "  Resumo: <uma linha>",
            "  Comportamento observado: <o que aconteceu>",
            "  Comportamento esperado: <referência ao caso de uso ou RN>",
            "  Reprodução: <passo a passo para o Claude reproduzir>",
            "  Evidência: <nome do screenshot, se houver>",
            "",
            "Agrupe por severidade. No final, gere um TL;DR com a contagem por sprint.'",
        ])

    make_content_slide(prs,
        "Passo 2 — Saída do ChatGPT (exemplo real)",
        "O que o ChatGPT devolve depois de processar suas observações",
        [
            "[UC-CV08] SEVERIDADE ALTA",
            "  Resumo: Decisão IA recomenda GO mesmo com score Jurídico abaixo do limiar",
            "  Observado: score_final=78, jur=30, IA retornou 'GO'",
            "  Esperado: RN-049 exige jur≥60 para GO; deveria recomendar 'Acompanhar'",
            "  Reprodução: login RP3X → Validação → edital hemograma → aba Aderência →",
            "    Calcular Scores → observar badge",
            "  Evidência: prints/sprint2/cv08_scores.png",
            "",
            "[UC-CV12] SEVERIDADE MÉDIA",
            "  Resumo: Gráfico de pizza da análise de mercado não renderiza",
            "  Observado: container em branco, sem mensagem de erro",
            "  Esperado: gráfico com 4 fatias mostrando distribuição por região",
            "  Reprodução: login RP3X → Validação → aba Análise de Mercado",
            "  Evidência: prints/sprint2/cv12.png",
        ])

    make_content_slide(prs,
        "Passo 3 — COLAR no documento de observações da sprint",
        "Crie UM arquivo por sprint com a saída do ChatGPT",
        [
            "Copie TUDO que o ChatGPT devolveu no Passo 2",
            "Cole num arquivo de texto — Google Docs, Notion, .md ou .txt",
            "Nome sugerido: observacoes_sprint[N]_[seu-nome].md",
            "  Exemplo: observacoes_sprint2_joao.md",
            "Adicione no topo do arquivo:",
            "  - Data da validação",
            "  - Empresa de teste usada (RP3X)",
            "  - Versão do sistema (consulte rodapé do Facilitia.ia)",
            "  - Total de UCs validados / Total de achados",
            "Anexe os screenshots numa pasta prints/sprintN/ ao lado do arquivo .md",
        ])

    make_content_slide(prs,
        "Passo 4 — ENTREGAR ao Claude Code",
        "Como enviar o documento de observações para correção",
        [
            "Abra o chat do Claude no claude.ai (ou no Claude Code dentro do projeto)",
            "Anexe o arquivo observacoes_sprint[N]_[nome].md (botão de clipe ou drag-drop)",
            "Anexe também os screenshots da pasta prints/sprintN/",
            "Mensagem ao Claude (texto-padrão):",
            "  'Validação manual da Sprint [N] concluída. Achados anexos.",
            "   Por favor: (1) priorize por severidade, (2) investigue cada item no código,",
            "   (3) corrija ou explique por que o comportamento atual é o correto,",
            "   (4) atualize os testes Playwright se necessário, (5) abra um commit por achado.'",
            "O Claude lê tudo, responde com plano de correção e executa sob sua revisão",
        ])

    make_content_slide(prs,
        "Exemplo MINI completo — uma observação até virar commit",
        "Vendo o ciclo inteiro funcionando para o caso do CRM",
        [
            "1) VOCÊ digita no ChatGPT (Passo 1):",
            "   'kanban CRM só mostra 8 colunas mas o UC-CRM01 fala em 13. Drag não funciona.'",
            "",
            "2) ChatGPT devolve (Passo 2):",
            "   '[UC-CRM01] CRÍTICA — Pipeline exibe 8 stages ao invés de 13 (RN-165).",
            "    Drag-and-drop inoperante. Reproduzir: login RP3X → CRM → aba Pipeline.'",
            "",
            "3) VOCÊ cola em observacoes_sprint5_joao.md e anexa no Claude (Passos 3-4)",
            "",
            "4) Claude Code:",
            "   - lê o arquivo, abre backend/crm_routes.py + frontend/CrmKanban.jsx",
            "   - identifica que faltam 5 stages no enum + DnD desabilitado por bug de id",
            "   - corrige, roda Playwright, gera commit 'fix(crm): 13 stages + drag&drop'",
            "   - responde com diff e link do commit pra você revisar",
        ])

    # ═══════════════════════════════════════════════════
    # SEÇÃO B.4: Fases Alfa e Beta Teste
    # ═══════════════════════════════════════════════════

    make_content_slide(prs,
        "Após o registro: Fase Alfa de Teste",
        "Teste interno fechado — primeira camada de qualidade",
        [
            "QUEM: validadores internos da empresa (time técnico + comercial)",
            "ESCOPO: todas as funcionalidades das 5 sprints em ambiente de homologação",
            "DURAÇÃO: 2 semanas, com rodadas diárias de feedback via ChatGPT→Claude",
            "OBJETIVO: caçar bugs críticos ANTES de expor a usuários externos",
            "SAÍDA: lista priorizada de correções + aprovação para avançar à fase Beta",
        ])

    make_content_slide(prs,
        "Fase Beta de Teste",
        "Teste aberto com usuários reais (empresas piloto)",
        [
            "QUEM: 3-5 empresas parceiras do setor de saúde (CH Hospitalar, RP3X, outras)",
            "ESCOPO: uso em cenários reais — editais verdadeiros do PNCP, propostas reais",
            "DURAÇÃO: 4-6 semanas, com suporte dedicado do time",
            "OBJETIVO: validar usabilidade, performance, precisão da IA em escala",
            "SAÍDA: relatório de maturidade + green-light para lançamento oficial",
        ])

    # ═══════════════════════════════════════════════════
    # SEÇÃO C: Sprint 1 (slides 13-21)
    # Screenshots: tutorial-sprint1-1/F##_*
    # ═══════════════════════════════════════════════════

    make_sprint_cover(prs, 1,
        "Construindo a Fundação",
        "17 casos de uso • RF-001 a RF-018",
        [
            "Cadastro completo da empresa com busca automática por CNPJ",
            "Portfolio de produtos com IA (consulta ANVISA, reprocessamento automático)",
            "Configurações: pesos GO/NO-GO, custos, regiões, fontes de busca",
        ])

    # Slide 14 — Cadastro empresa (screenshot: F01_06_salvo = dados salvos da CH Hospitalar)
    make_content_slide(prs,
        "Cadastro da empresa",
        "UC-F01, UC-F02 | RF-001, RF-005",
        [
            "O primeiro passo: cadastrar CNPJ, razão social, endereço, contatos",
            "O sistema busca dados automaticamente pelo CNPJ",
            "Na screenshot: dados da CH Hospitalar salvos com sucesso",
            "Inclui e-mails, telefones e área de atuação padrão (UC-F02)",
        ],
        screenshot_uc="UC-F01",
        footer="UC-F01/F02 | RF-001, RF-005 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md | Validação: RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        sprint_num=1)

    # Slide 15 — Documentos e certidões (screenshot: F03_05 = lista de docs + F04_11 = certidões OK)
    make_content_slide(prs,
        "Documentos e certidões",
        "UC-F03, UC-F04 | RF-002, RF-004",
        [
            "Upload de documentos com tipo, validade e PDF anexo",
            "Busca automática de certidões (CND, FGTS) com frequência configurável",
            "Na screenshot: lista de documentos cadastrados para a CH Hospitalar",
            "O sistema alerta quando uma certidão está para vencer",
        ],
        screenshot_uc="UC-F03",
        footer="UC-F03/F04 | RF-002, RF-004 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md | Validação: RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        sprint_num=1)

    # Slide 16 — Responsáveis técnicos (screenshot: F05_05 = lista de 3 responsáveis)
    make_content_slide(prs,
        "Responsáveis técnicos",
        "UC-F05 | RF-003",
        [
            "Cadastro de responsáveis técnicos e administrativos",
            "Nome, cargo, CPF, registro profissional (CRM, CREA etc.)",
            "Na screenshot: lista de responsáveis cadastrados para CH Hospitalar",
            "A IA consulta estes dados ao montar propostas e petições",
        ],
        screenshot_uc="UC-F05",
        footer="UC-F05 | RF-003 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md | Validação: RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        sprint_num=1)

    # Slide 17 — Portfolio de produtos (screenshot: F06_04 = detalhe de produto)
    make_content_slide(prs,
        "Portfolio de produtos",
        "UC-F06, UC-F07 | RF-006 a RF-012",
        [
            "O coração do sistema: os produtos que a empresa vende",
            "Cadastro manual, via CSV, ou pela IA (consulta ANVISA + manual técnico)",
            "Na screenshot: tela de portfólio com listagem de produtos e filtros",
            "Listagem com filtros por área, classe, status e busca textual",
        ],
        screenshot_uc="UC-F06",
        footer="UC-F06/F07 | RF-006 a RF-012 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md | Validação: RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        sprint_num=1)

    # Slide 18 — Edição e IA nos produtos
    make_content_slide(prs,
        "Edição e IA nos produtos",
        "UC-F08 a UC-F13 | RF-008, RF-010, RF-012",
        [
            "Edição de especificações técnicas e NCM (UC-F08)",
            "IA reprocessa specs a partir de manuais técnicos (UC-F09)",
            "Consulta ANVISA e busca web para dados complementares (UC-F10)",
            "Verificação de completude com scores (UC-F11)",
            "Na screenshot: produto reprocessado pela IA com specs atualizadas",
        ],
        screenshot_uc="UC-F09",
        footer="UC-F08..F13 | RF-008, RF-010, RF-012 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md",
        sprint_num=1)

    # Slide 19 — Pesos GO/NO-GO (screenshot: F14_06 = limiares salvos)
    make_content_slide(prs,
        "Pesos GO/NO-GO",
        "UC-F14 | RF-018",
        [
            "Como a empresa define o que é mais importante para ela?",
            "6 dimensões com pesos: Técnico, Comercial, Jurídico, Logístico, Documental, Complexidade",
            "Limiares ajustáveis: abaixo de X = NO-GO automático",
            "Na screenshot: limiares configurados e salvos para CH Hospitalar",
        ],
        screenshot_uc="UC-F14",
        footer="UC-F14 | RF-018 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md | Validação: RELATORIO_VALIDACAO_SPRINT1_TUTORIAL1.md",
        sprint_num=1)

    # Slide 20 — Parametrizações gerais
    make_content_slide(prs,
        "Parametrizações gerais",
        "UC-F15 a UC-F17 | RF-013 a RF-018",
        [
            "Custos fixos mensais, markup padrão, frete base (UC-F15)",
            "Fontes de busca: PNCP, Brave, Google, BEC-SP (UC-F16)",
            "Palavras-chave e NCMs de interesse para captação automática",
            "Notificações e preferências do sistema (UC-F17)",
            "Na screenshot: NCMs de interesse configurados e salvos",
        ],
        screenshot_uc="UC-F16",
        footer="UC-F15..F17 | RF-013 a RF-018 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md",
        sprint_num=1)

    # Slide novo V3 — Spotlight Sprint 1: Formulário empresa (UC-F01)
    _set_section_kicker("SPRINT 1 · SPOTLIGHT")
    content_slide_with_excerpt(prs,
        kicker="SPRINT 1 · SPOTLIGHT",
        title="Cadastro de empresa — o que observar",
        subtitle="UC-F01 · dados carregados automaticamente pelo CNPJ",
        bullets=[
            "Campo CNPJ dispara lookup automático na ReceitaWS",
            "Razão social, endereço e CNAE preenchidos sem digitação",
            "Validação DV CNPJ antes de salvar (RN-028)",
            "Botão Salvar ativo só quando campos obrigatórios preenchidos",
        ],
        screenshot_uc="UC-F01",
        excerpt_key="uc_sprint1",
        callouts_rel=[
            (0.05, 0.10, 0.60, 0.15, "Campo CNPJ + lookup"),
            (0.05, 0.28, 0.90, 0.15, "Razão social autopreenchida"),
            (0.05, 0.45, 0.90, 0.40, "Endereço e contatos"),
        ],
        footer="UC-F01 | RF-001 | Doc: CASOS DE USO EMPRESA PORTFOLIO PARAMETRIZACAO V2.md",
        sprint_num=1)

    # ═══════════════════════════════════════════════════
    # SEÇÃO D: Sprint 2 (slides 22-30)
    # Screenshots: validacao-sprint2/CV##_*
    # ═══════════════════════════════════════════════════

    make_sprint_cover(prs, 2,
        "Encontrar e Avaliar Oportunidades",
        "13 casos de uso • RF-019 a RF-037",
        [
            "Busca automática de editais no PNCP (Portal Nacional de Contratações)",
            "Scores em 6 dimensões com recomendação GO/NO-GO da IA",
            "Monitoramentos automáticos para novos editais sem buscar todo dia",
        ])

    # Slide 22 — Buscar editais (screenshot: CV01-P05_resp = resultados de busca com filtros)
    make_content_slide(prs,
        "Buscar editais no PNCP",
        "UC-CV01 | RF-019, RF-021, RF-022, RF-026",
        [
            "O usuário digita um termo (ex: 'hemograma') e aplica filtros",
            "O sistema consulta o PNCP e retorna editais com scores automáticos",
            "Na screenshot: tela de Captação com resultados de busca e filtros aplicados",
            "5 cenários de busca testados na validação (monitor, ultrassom, etc.)",
        ],
        screenshot_uc="UC-CV01",
        footer="UC-CV01 | RF-019, RF-021, RF-022 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 23 — Painel de detalhes (screenshot: CV02_resp_painel)
    make_content_slide(prs,
        "Painel de detalhes do edital",
        "UC-CV02 | RF-019, RF-020, RF-024",
        [
            "Ao clicar num edital, abre um painel lateral deslizante",
            "Visão completa: scores detalhados, itens, órgão, valor, prazo",
            "Na screenshot: painel lateral com detalhes do edital selecionado",
            "É a visão que o analista usa para decidir antes do GO/NO-GO",
        ],
        screenshot_uc="UC-CV02",
        footer="UC-CV02 | RF-019, RF-020, RF-024 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 24 — Salvar e estratégia (screenshot: CV03-P02_resp_salvos)
    make_content_slide(prs,
        "Salvar edital e definir estratégia",
        "UC-CV03, UC-CV04 | RF-019, RF-023, RF-027",
        [
            "Gostou do edital? Salve no banco com itens e scores",
            "Na screenshot: tela de resultados com editais após salvamento",
            "Defina estratégia: participar, observar ou declinar (UC-CV04)",
            "Configure margem comercial e intenção específica",
        ],
        screenshot_uc="UC-CV03",
        footer="UC-CV03/CV04 | RF-019, RF-023, RF-027 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 25 — Exportar e monitorar (screenshot: CV05_resp_exportado)
    make_content_slide(prs,
        "Exportar e monitorar automaticamente",
        "UC-CV05, UC-CV06 | RF-019, RF-025, RF-026",
        [
            "Exporte resultados em CSV/PDF para compartilhar com a equipe",
            "Na screenshot: resultado de busca exportado com sucesso",
            "Configure monitoramentos periódicos para receber alertas (UC-CV06)",
            "Novos editais são detectados sem buscar manualmente",
        ],
        screenshot_uc="UC-CV05",
        footer="UC-CV05/CV06 | RF-019, RF-025, RF-026 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 26 — Scores e GO/NO-GO (screenshot: CV08_resp_scores)
    make_content_slide(prs,
        "Scores e decisão GO/NO-GO",
        "UC-CV07, UC-CV08 | RF-027, RF-028, RF-037",
        [
            "O momento crucial: a IA calcula 6 scores e recomenda GO ou NO-GO",
            "Técnico, Comercial, Jurídico, Logístico, Documental, Complexidade",
            "Na screenshot: scores calculados para o edital com cores e valores",
            "Justificativa escrita pela IA explicando a recomendação",
        ],
        screenshot_uc="UC-CV08",
        footer="UC-CV07/CV08 | RF-027, RF-028, RF-037 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 27 — Importação de itens (screenshot: CV09_resp_itens)
    make_content_slide(prs,
        "Importação de itens e lotes",
        "UC-CV09 | RF-031, RF-036",
        [
            "O sistema importa itens do edital diretamente do PNCP",
            "Agrupamento automático em lotes via IA (por similaridade)",
            "Na screenshot: itens do edital importados na aba de lotes",
            "Cada item traz: descrição, quantidade, unidade, valor estimado",
        ],
        screenshot_uc="UC-CV09",
        footer="UC-CV09 | RF-031, RF-036 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 28 — Análises avançadas (screenshot: CV11_resp_riscos)
    make_content_slide(prs,
        "Análises avançadas",
        "UC-CV10 a UC-CV12 | RF-029 a RF-034",
        [
            "Confronto automático: documentos exigidos vs documentos da empresa (UC-CV10)",
            "Na screenshot: análise de riscos com classificação de gravidade (UC-CV11)",
            "Recorrência do órgão, atas anteriores e concorrentes (UC-CV11)",
            "Análise de mercado: histórico de compras e volume (UC-CV12)",
        ],
        screenshot_uc="UC-CV11",
        footer="UC-CV10..CV12 | RF-029 a RF-034 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide 29 — Chat com IA (screenshot: CV13_resp_resposta)
    make_content_slide(prs,
        "Chat com IA na validação",
        "UC-CV13 | RF-026, RF-029, RF-030",
        [
            "Converse com a IA sobre o edital em linguagem natural",
            "'Resuma este edital', 'Quais os riscos?', 'Recomende ação'",
            "Na screenshot: resposta da IA com análise detalhada do edital",
            "Respostas baseadas nos dados reais do edital e da empresa",
        ],
        screenshot_uc="UC-CV13",
        footer="UC-CV13 | RF-026, RF-029, RF-030 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md | Validação: RESULTADO VALIDACAO SPRINT2.md",
        sprint_num=2)

    # Slide novo V3 — Spotlight Sprint 2: ScoreBars 6D (UC-CV08)
    _set_section_kicker("SPRINT 2 · SPOTLIGHT")
    content_slide_with_excerpt(prs,
        kicker="SPRINT 2 · SPOTLIGHT",
        title="Scores 6D — o que observar",
        subtitle="UC-CV08 · decisão GO/NO-GO com recomendação IA",
        bullets=[
            "6 ScoreBars coloridas: técnico, financeiro, risco, margem, exclusividade, diferencial",
            "Agregado ponderado visível no topo (0 a 100)",
            "Recomendação GO/NO-GO com justificativa da IA",
            "Badge 'decidido em HH:MM' registra quando validador confirmou",
        ],
        screenshot_uc="UC-CV08",
        excerpt_key="uc_sprint2",
        callouts_rel=[
            (0.05, 0.10, 0.90, 0.10, "Score agregado 0-100 + badge decisão"),
            (0.05, 0.25, 0.90, 0.50, "6 ScoreBars dimensionais"),
            (0.05, 0.78, 0.90, 0.18, "Recomendação IA + justificativa"),
        ],
        footer="UC-CV08 | RF-027, RF-028, RF-037 | Doc: CASOS DE USO CAPTACAO VALIDACAO(SPRINT2) V2.md",
        sprint_num=2)

    # ═══════════════════════════════════════════════════
    # SEÇÃO E: Sprint 3 (slides 31-40)
    # Screenshots: validacao-sprint3/P##_* e R##_*
    # ═══════════════════════════════════════════════════

    make_sprint_cover(prs, 3,
        "Precificação em 6 Camadas + Proposta Automática",
        "19 casos de uso • RF-039 a RF-041",
        [
            "6 camadas de precificação: Volumetria → Custos → Preço Base → Referência → Lances → Histórico",
            "Motor de proposta técnica gerado automaticamente pela IA",
            "Auditoria ANVISA + checklist documental + exportação de dossiê",
        ])

    # Slide 31 — Organização por lotes (screenshot: P01_p7_tabela_itens)
    make_content_slide(prs,
        "Organização por lotes",
        "UC-P01 | RF-039-01",
        [
            "Os itens do edital são organizados em lotes com parâmetros técnicos",
            "Cada lote pode ter sua própria estratégia de preço",
            "Na screenshot: tabela de itens do edital organizados em lotes",
        ],
        screenshot_uc="UC-P01",
        footer="UC-P01 | RF-039-01 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 32 — Seleção + Volumetria
    make_content_slide(prs,
        "Seleção inteligente + Volumetria",
        "UC-P02, UC-P03 | RF-039-02, RF-039-07",
        [
            "A IA sugere quais produtos do portfolio atendem cada item do edital",
            "Na screenshot: lista de lotes com status de vinculação (UC-P02)",
            "Cálculo automático de quantidades técnicas por parâmetro (UC-P03)",
        ],
        screenshot_uc="UC-P02",
        footer="UC-P02/P03 | RF-039-02, RF-039-07 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 33 — Custos ERP
    make_content_slide(prs,
        "Custos ERP e tributário",
        "UC-P04 | RF-039-03, RF-039-04",
        [
            "Base de custos integrada: custo do produto + frete + impostos",
            "Regras tributárias por NCM aplicadas automaticamente",
            "Na screenshot: custos salvos com valores e impostos calculados",
        ],
        screenshot_uc="UC-P04",
        footer="UC-P04 | RF-039-03, RF-039-04 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 34 — Preço base e referência
    make_content_slide(prs,
        "Preço base e referência — Camadas B e C",
        "UC-P05, UC-P06 | RF-039-08, RF-039-09",
        [
            "Camada B: 3 modos — manual, custos + markup, ou histórico",
            "Na screenshot: preço base definido e salvo (UC-P05)",
            "Camada C: valor de referência do edital para comparação (UC-P06)",
        ],
        screenshot_uc="UC-P05",
        footer="UC-P05/P06 | RF-039-08, RF-039-09 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 35 — Lances e estratégia
    make_content_slide(prs,
        "Lances e estratégia competitiva",
        "UC-P07, UC-P08 | RF-039-10, RF-039-11",
        [
            "Camada D: primeiro lance, Camada E: lance mínimo",
            "Na screenshot: lances salvos com valores definidos (UC-P07)",
            "Estratégia agressiva ou conservadora com simulador de disputa (UC-P08)",
        ],
        screenshot_uc="UC-P07",
        footer="UC-P07/P08 | RF-039-10, RF-039-11 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 36 — Histórico + Comodato
    make_content_slide(prs,
        "Histórico de preços + Comodato",
        "UC-P09, UC-P10 | RF-039-12, RF-039-13",
        [
            "Camada F: preços de editais similares dos últimos 2 anos",
            "Na screenshot: resultado de busca histórica com filtros (UC-P09)",
            "Gestão de equipamentos em comodato com prazos (UC-P10)",
        ],
        screenshot_uc="UC-P09",
        footer="UC-P09/P10 | RF-039-12, RF-039-13 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 37 — Pipeline IA + Relatório
    make_content_slide(prs,
        "Pipeline IA + Relatório consolidado",
        "UC-P11, UC-P12 | RF-039-14, RF-039-15",
        [
            "A IA analisa TUDO: as 6 camadas + portfolio + edital",
            "Na screenshot: confirmação de precificação analisada pela IA (UC-P11)",
            "Relatório consolidado em HTML com recomendação de preço (UC-P12)",
        ],
        screenshot_uc="UC-P11",
        footer="UC-P11/P12 | RF-039-14, RF-039-15 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 38 — Motor de proposta técnica
    make_content_slide(prs,
        "Motor de proposta técnica",
        "UC-R01 a UC-R05 | RF-040",
        [
            "O sistema gera a proposta técnica automaticamente com dados do edital",
            "Na screenshot: página de detalhes da proposta com itens e parecer (UC-R01)",
            "Upload de proposta externa (UC-R02) + edição A/B (UC-R03)",
            "Auditoria ANVISA com semáforo regulatório (UC-R04) + checklist documental (UC-R05)",
        ],
        screenshot_uc="UC-R01",
        footer="UC-R01..R05 | RF-040 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide 39 — Exportação e submissão
    make_content_slide(prs,
        "Exportação e submissão",
        "UC-R06, UC-R07 | RF-041",
        [
            "Dossiê completo exportado em PDF, DOCX ou ZIP",
            "Na screenshot: seleção de documentos/análises para exportação com checkboxes (UC-R06)",
            "Status rastreado: rascunho → enviada → aceita (UC-R07)",
        ],
        screenshot_uc="UC-R06",
        footer="UC-R06/R07 | RF-041 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md | Validação: RESULTADO VALIDACAO SPRINT3.md",
        sprint_num=3)

    # Slide novo V3 — Spotlight Sprint 3: Simulador de disputa (UC-P08)
    _set_section_kicker("SPRINT 3 · SPOTLIGHT")
    content_slide_with_excerpt(prs,
        kicker="SPRINT 3 · SPOTLIGHT",
        title="Simulador de disputa — o que observar",
        subtitle="UC-P08 · Camada E — lance mínimo e margem",
        bullets=[
            "Input do lance do concorrente estimado",
            "Slider de decremento percentual simulando a disputa",
            "Cálculo em tempo real do lucro projetado",
            "Bloqueio em vermelho quando margem < mínima (RN-039-10)",
        ],
        screenshot_uc="UC-P08",
        excerpt_key="resultado_sprint3",
        callouts_rel=[
            (0.05, 0.15, 0.40, 0.25, "Input lance do concorrente"),
            (0.05, 0.45, 0.40, 0.20, "Slider decremento %"),
            (0.50, 0.15, 0.45, 0.60, "Painel resultado: lance mínimo + margem"),
        ],
        footer="UC-P08 | RF-039-10, RF-039-11 | Doc: CASOS DE USO PRECIFICACAO E PROPOSTA V2.md",
        sprint_num=3)

    # ═══════════════════════════════════════════════════
    # SEÇÃO F: Sprint 4 (slides 41-47)
    # Screenshots: validacao-sprint4/I##_*, RE##_*, FU##_*
    # ═══════════════════════════════════════════════════

    make_sprint_cover(prs, 4,
        "Defesa Jurídica Assistida por IA",
        "11 casos de uso • RF-043, RF-044 • Lei 14.133/2021",
        [
            "Análise de conformidade legal com classificação de gravidade",
            "Geração de petições de impugnação com fundamentação jurídica pela IA",
            "Recursos e contra-razões com motivações estruturadas",
        ])

    # Slide 41 — Validação legal (screenshot: I01_resp_tabela_inconsistencias)
    make_content_slide(prs,
        "Validação legal do edital",
        "UC-I01, UC-I02 | RF-043-01, RF-043-02",
        [
            "A IA analisa o edital contra a Lei 14.133/2021 e regulamentos",
            "Classifica inconsistências: ALTA, MÉDIA ou BAIXA gravidade",
            "Na screenshot: tabela de inconsistências com artigos de lei e gravidade",
            "Sugere esclarecimento ou impugnação para cada achado (UC-I02)",
        ],
        screenshot_uc="UC-I01",
        footer="UC-I01/I02 | RF-043-01, RF-043-02 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide 42 — Petição de impugnação (screenshot: I03_resp_peticao_criada)
    make_content_slide(prs,
        "Petição de impugnação",
        "UC-I03, UC-I04 | RF-043-03, RF-043-07",
        [
            "A IA gera minuta com fundamentação legal completa",
            "Na screenshot: modal de criação de petição com campos de artigos e argumentos",
            "Upload de petição redigida externamente (UC-I04)",
            "Edição completa com LOG de alterações",
        ],
        screenshot_uc="UC-I03",
        footer="UC-I03/I04 | RF-043-03, RF-043-07 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide 43 — Controle de prazo (screenshot: I05_resp_countdown)
    make_content_slide(prs,
        "Controle de prazo",
        "UC-I05 | RF-043-08",
        [
            "Countdown visual: quantos dias faltam para impugnar",
            "Na screenshot: tabela de prazos de impugnação com status e prazo restante",
            "Cores: verde (>5 dias), amarelo (3-5 dias), vermelho (<3 dias)",
            "Alertas automáticos quando o prazo é crítico",
        ],
        screenshot_uc="UC-I05",
        footer="UC-I05 | RF-043-08 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide 44 — Recursos (screenshot: RE01_resp_timeline)
    make_content_slide(prs,
        "Recursos — monitoramento e análise",
        "UC-RE01 a UC-RE03 | RF-044-01 a RF-044-03",
        [
            "Monitoramento da janela de recurso com timeline visual (UC-RE01)",
            "Na screenshot: badge de monitoramento ativo com acesso aos detalhes",
            "Análise da proposta vencedora para identificar fragilidades (UC-RE02)",
            "Chat jurídico interativo para perguntas específicas (UC-RE03)",
        ],
        screenshot_uc="UC-RE01",
        footer="UC-RE01..RE03 | RF-044-01 a RF-044-03 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide 45 — Laudos (screenshot: RE04_resp_laudo_criado)
    make_content_slide(prs,
        "Laudos de recurso e contra-razão",
        "UC-RE04, UC-RE05 | RF-044-07, RF-044-08",
        [
            "Recurso: quando a empresa perdeu e quer contestar (UC-RE04)",
            "Na screenshot: laudo de recurso criado com motivações estruturadas",
            "Contra-razão: quando a empresa ganhou e precisa se defender (UC-RE05)",
            "Geração automática com categorias e argumentos pela IA",
        ],
        screenshot_uc="UC-RE04",
        footer="UC-RE04/RE05 | RF-044-07, RF-044-08 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide 46 — Submissão assistida (screenshot: RE06_resp_submissao_registrada)
    make_content_slide(prs,
        "Submissão assistida no portal",
        "UC-RE06 | RF-044-12",
        [
            "O sistema auxilia o preenchimento nos portais da administração",
            "Na screenshot: submissão registrada com protocolo",
            "Dados pré-preenchidos a partir do laudo gerado",
        ],
        screenshot_uc="UC-RE06",
        footer="UC-RE06 | RF-044-12 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=4)

    # Slide novo V3 — Spotlight Sprint 4: Tabela de inconsistências legais
    _set_section_kicker("SPRINT 4 · SPOTLIGHT")
    content_slide_with_excerpt(prs,
        kicker="SPRINT 4 · SPOTLIGHT",
        title="Tabela de inconsistências — o que observar",
        subtitle="UC-I01 · análise legal do edital pela IA",
        bullets=[
            "Coluna 'Gravidade' com selo ALTA/MÉDIA/BAIXA colorido",
            "Artigo da Lei 14.133/2021 citado em cada linha",
            "Sugestão de ação: esclarecimento ou impugnação",
            "A IA justifica cada achado com citação legal",
        ],
        screenshot_uc="UC-I01",
        excerpt_key="resultado_sprint4",
        callouts_rel=[
            (0.02, 0.18, 0.22, 0.70, "Coluna Artigo/Lei"),
            (0.60, 0.18, 0.22, 0.70, "Coluna Gravidade (selos)"),
            (0.82, 0.18, 0.16, 0.70, "Coluna Ação sugerida"),
        ],
        footer="UC-I01 | RF-043-01 | Doc: CASOS DE USO RECURSOS E IMPUGNACOES V2.md",
        sprint_num=4)

    # ═══════════════════════════════════════════════════
    # SEÇÃO G: Sprint 5 (slides 48-63)
    # Screenshots: UC-CT##/ e UC-CRM##/ (P##_resp.png)
    # ═══════════════════════════════════════════════════

    make_sprint_cover(prs, 5,
        "Pós-venda: Execução, Contratos e CRM",
        "19 casos de uso • RF-017, RF-035, RF-045, RF-046, RF-051, RF-052",
        [
            "Follow-up de resultados, atas de pregão, gestão de contratos",
            "Empenhos encadeados com auditoria cruzada automática",
            "CRM completo: pipeline 13 stages, mapa, agenda, KPIs, decisões",
        ])

    # Slide 48 — Follow-up
    make_content_slide(prs,
        "Follow-up de resultados",
        "UC-FU01 a UC-FU03 | RF-017, RF-046",
        [
            "Registrar resultado: vitória, derrota ou cancelamento (UC-FU01)",
            "Na screenshot: página de follow-up com contadores (1 ganho, 0 parcial) e editais pendentes",
            "Configurar alertas para marcos da execução (UC-FU02)",
            "Score logístico de desempenho calculado automaticamente (UC-FU03)",
        ],
        screenshot_uc="UC-FU01",
        footer="UC-FU01..FU03 | RF-017, RF-046 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT4.md",
        sprint_num=5)

    # Slide 49 — Atas de pregão
    make_content_slide(prs,
        "Atas de pregão",
        "UC-AT01 a UC-AT03 | RF-035",
        [
            "Buscar atas diretamente no PNCP pelo número do processo",
            "Extrair vencedor, valor adjudicado e detalhes da ata em PDF",
            "Dashboard com histórico de atas consultadas e filtros",
        ],
        screenshot_uc="UC-AT01",
        footer="UC-AT01..AT03 | RF-035 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 50 — Contratos cadastro
    make_content_slide(prs,
        "Contratos — cadastro e entregas",
        "UC-CT01 a UC-CT03 | RF-046",
        [
            "Cadastrar contrato com dados: número, valor, vigência, objeto",
            "Na screenshot: listagem de contratos com stat cards (12 contratos, R$ 5.950.000)",
            "Registrar entregas com número NF, data e status (UC-CT02)",
            "Cronograma previsto × realizado com alertas de atraso (UC-CT03)",
        ],
        screenshot_uc="UC-CT01",
        footer="UC-CT01..CT03 | RF-046 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 51 — Aditivos
    make_content_slide(prs,
        "Aditivos, designações e ARP",
        "UC-CT04 a UC-CT06 | RF-046",
        [
            "Gestão de aditivos contratuais (prazo, valor, objeto)",
            "Designar gestor e fiscal do contrato com portaria",
            "Controle de saldo ARP (autorizado não recebido) e caronas",
            "Na screenshot: aba Aditivos na página de Execução de Contratos",
        ],
        screenshot_uc="UC-CT04",
        footer="UC-CT04..CT06 | RF-046 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 52 — Empenhos (screenshot: UC-CT07/P03_resp = stat cards + tabela)
    make_content_slide(prs,
        "Gestão de Empenhos",
        "UC-CT07 | RF-046-01",
        [
            "Tabela de empenhos: número EMPH-2026, tipo, valor, data",
            "Stat cards: total empenhado, faturado, saldo (com % colorido)",
            "Na screenshot: stat cards com valores e tabela de empenhos",
            "Itens incluem calibradores SEM VALOR com alerta visual",
            "Botão Novo Empenho para criar empenhos adicionais",
        ],
        screenshot_uc="UC-CT07",
        footer="UC-CT07 | RF-046-01 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 53 — Auditoria (screenshot: UC-CT08/P03_resp = cruzamento)
    make_content_slide(prs,
        "Auditoria Empenho × Fatura × Entrega",
        "UC-CT08 | RF-046-02",
        [
            "Cruzamento automático: empenhado vs faturado vs pago vs entregue",
            "Na screenshot: totais coloridos e tabela com divergências detectadas",
            "Detecta divergências com alerta laranja (ex: R$24.000 de diferença)",
            "Exportar relatório em CSV para auditoria externa",
        ],
        screenshot_uc="UC-CT08",
        footer="UC-CT08 | RF-046-02 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 54 — Contratos a vencer (screenshot: UC-CT09/P03_resp = 5 tiers)
    make_content_slide(prs,
        "Contratos a Vencer — 5 tiers",
        "UC-CT09 | RF-046-03",
        [
            "Cards organizados em 5 faixas de vencimento:",
            "Vence em 30 dias (urgente) — Vence em 90 dias (atenção)",
            "Em tratativa — Renovado — Não renovado",
            "Na screenshot: grid de cards coloridos em 5 colunas (tiers)",
        ],
        screenshot_uc="UC-CT09",
        footer="UC-CT09 | RF-046-03 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 55 — KPIs Execução (screenshot: UC-CT10/P03_resp = stat cards)
    make_content_slide(prs,
        "KPIs de Execução",
        "UC-CT10 | RF-046-04",
        [
            "6 stat cards numéricos com dados reais",
            "Contratos Ativos, Vence 30d, Vence 90d, Em Tratativa, Renovados, Não Renovados",
            "Na screenshot: indicadores com ícones coloridos e números",
            "Filtro de período para recálculo dinâmico",
        ],
        screenshot_uc="UC-CT10",
        footer="UC-CT10 | RF-046-04 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 56 — Contratado × Realizado
    make_content_slide(prs,
        "Contratado × Realizado",
        "UC-CR01 a UC-CR03 | RF-051, RF-052",
        [
            "Dashboard visual: valores contratados vs efetivamente recebidos",
            "Na screenshot: cards de totais, pedidos em atraso e próximos vencimentos",
            "Alertas multi-tier por gravidade com cores (rosa, amarelo, verde)",
        ],
        screenshot_uc="UC-CR01-acao",
        footer="UC-CR01..CR03 | RF-051, RF-052 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 57 — Pipeline CRM (screenshot: UC-CRM01/P03_resp = kanban)
    make_content_slide(prs,
        "Pipeline CRM — 13 stages",
        "UC-CRM01 | RF-045, RF-045-01",
        [
            "O CORAÇÃO do CRM: kanban horizontal com 13 colunas",
            "De 'Captado Não Divulgado' até 'Resultados Definitivos'",
            "Na screenshot: grid de cards organizados por stage",
            "Cada card mostra: edital, órgão, valor estimado, modalidade",
            "151 editais distribuídos no pipeline (dados reais do seed)",
        ],
        screenshot_uc="UC-CRM01",
        footer="UC-CRM01 | RF-045-01 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 58 — Parametrizações CRM (screenshot: UC-CRM02/P03_resp = tabela)
    make_content_slide(prs,
        "Parametrizações CRM",
        "UC-CRM02 | RF-045-02",
        [
            "3 sub-abas configuráveis pela empresa:",
            "Tipos de Edital (8 padrões), Agrupamentos Portfolio (12), Motivos de Derrota (7)",
            "Na screenshot: tabela de parametrizações com status (ativo/inativo)",
            "A empresa personaliza conforme sua realidade de mercado",
        ],
        screenshot_uc="UC-CRM02",
        footer="UC-CRM02 | RF-045-02 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide — Mapa geográfico
    make_content_slide(prs,
        "Mapa geográfico de editais",
        "UC-CRM03 | RF-045-03",
        [
            "Distribuição geográfica dos 151 editais captados por UF",
            "Cards por estado com contagem e valor total",
            "Na screenshot: SP 20, GO 11, PR 10, RS 9, SC 5, MG/BA/CE/RJ/PA/TO menores",
            "Clicar num estado filtra o pipeline pela UF selecionada",
        ],
        screenshot_uc="UC-CRM03",
        footer="UC-CRM03 | RF-045-03 | Doc: CASOS DE USO SPRINT5 V4.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide — Agenda de compromissos
    make_content_slide(prs,
        "Agenda de compromissos",
        "UC-CRM04 | RF-045-04",
        [
            "Lista de compromissos vinculados aos editais em andamento",
            "Badges de urgência: CRÍTICA (vermelho), ALTA (laranja), NORMAL (azul)",
            "Na screenshot: 12 itens com data de abertura do pregão e órgão",
            "Permite agendar visitas, reuniões e entregas de amostras",
        ],
        screenshot_uc="UC-CRM04",
        footer="UC-CRM04 | RF-045-04 | Doc: CASOS DE USO SPRINT5 V4.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 60 — KPIs CRM (screenshot: UC-CRM05/P03_resp = stat cards)
    make_content_slide(prs,
        "KPIs CRM",
        "UC-CRM05 | RF-045-05",
        [
            "8 stat cards com números reais:",
            "Total (101), Leads (29), Propostas (20), Perdidos (3)",
            "Na screenshot: métricas em cards com ícones coloridos",
            "Taxas de conversão e ticket médio calculados automaticamente",
        ],
        screenshot_uc="UC-CRM05",
        footer="UC-CRM05 | RF-045-05 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide 61 — Decisões (screenshot: UC-CRM06/P03_resp = tabela de decisões)
    make_content_slide(prs,
        "Decisões: não-participação e perda",
        "UC-CRM06, UC-CRM07 | RF-045-01",
        [
            "Registro estruturado de: por que NÃO participou e por que PERDEU",
            "Na screenshot: tabela de decisões com motivos, datas e status",
            "Não-participação: motivo (ex: exclusivo ME/EPP) + justificativa",
            "Perda: categoria + descrição + checkbox contra-razão",
        ],
        screenshot_uc="UC-CRM06",
        footer="UC-CRM06/CRM07 | RF-045-01 | Doc: CASOS DE USO SPRINT5 V3.md | Validação: RESULTADO VALIDACAO SPRINT5.md",
        sprint_num=5)

    # Slide novo V3 — Spotlight Sprint 5: Pipeline CRM com callouts
    _set_section_kicker("SPRINT 5 · SPOTLIGHT")
    content_slide_with_excerpt(prs,
        kicker="SPRINT 5 · SPOTLIGHT",
        title="Pipeline CRM — o que observar",
        subtitle="UC-CRM01 · tela kanban com 13 colunas da Sprint 5",
        bullets=[
            "13 colunas de stages — da captação à pós-venda",
            "Drag-and-drop entre colunas muda o stage (RF-045-01)",
            "Contadores no topo de cada coluna mostram volume",
            "Cards coloridos por prioridade da IA",
        ],
        screenshot_uc="UC-CRM01",
        excerpt_key="resultado_sprint5",
        callouts_rel=[
            (0.02, 0.05, 0.96, 0.12, "Header com 13 stages visíveis"),
            (0.02, 0.20, 0.35, 0.60, "Primeira coluna (novos editais)"),
            (0.40, 0.20, 0.25, 0.60, "Coluna de análise/GO"),
        ],
        footer="UC-CRM01 | RF-045-01 | Doc: CASOS DE USO SPRINT5 V3.md",
        sprint_num=5)

    _set_section_kicker(None)  # reset antes do encerramento

    # ═══════════════════════════════════════════════════
    # SEÇÃO H: Encerramento (slides 62-66)
    # ═══════════════════════════════════════════════════

    # Slide novo V3 — Arquitetura de dados + IA
    diagram_slide(prs,
        kicker="ARQUITETURA · DADOS + IA",
        title="Fontes externas, backend, frontend e validador",
        subtitle="PNCP + DeepSeek + Brave + ANVISA → Flask → React → validador humano",
        draw_fn=draw_pipeline_ia_dados,
        footer="Fontes de dados reais · IA via OpenRouter · interface única pro validador")

    make_content_slide(prs,
        "Resumo de conformidade",
        "Todas as 5 sprints aprovadas com 100% de testes passando",
        [
            "Sprint 1: 17 UCs — Empresa, Portfolio, Parametrizações ✅",
            "Sprint 2: 13 UCs — Captação e Validação com IA ✅",
            "Sprint 3: 19 UCs — Precificação 6 Camadas + Proposta ✅",
            "Sprint 4: 11 UCs — Impugnação, Recursos, Contra-Razões ✅",
            "Sprint 5: 19 UCs — Follow-up, Atas, Contratos, CRM ✅",
        ])

    make_content_slide(prs,
        "O que vem pela frente",
        "Sprints 6 a 10 planejadas",
        [
            "Sprint 6: CRM leads + Dashboard de Perdas + Análise de Concorrência",
            "Sprint 7: Flags, Monitoria, Auditoria completa, Notificações SMTP",
            "Sprint 8: Mercado TAM/SAM/SOM + Analytics com MindsDB",
            "Sprint 9: Dispensas de licitação + Classes avançadas + Máscaras",
            "Sprint 10: Disputa de Lances + Health Check + QA end-to-end",
        ])

    make_content_slide(prs,
        "Como acessar para validar",
        "Tudo pronto para você testar",
        [
            "URL: http://192.168.1.115:5179 (instância de validação)",
            "Usuário CH Hospitalar: valida1@valida.com.br / 123456",
            "Usuário RP3X: valida2@valida.com.br / 123456 (use este para validação manual)",
            "Documentos: pasta docs/ (tutoriais-2, casos de uso, resultados de validação)",
            "Siga os tutorialsprint#-2.md na ordem: Sprint 1 → 2 → 3 → 4 → 5",
        ])

    # Slide final — Obrigado
    slide = new_slide(prs)
    add_band(slide, Inches(0), Inches(0), Inches(0.3), SLIDE_H, EMERALD)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
                "FACILITIA.IA", size=12, color=EMERALD, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(12), Inches(1.4),
                "Obrigado", size=72, color=TXT_INK, bold=True)
    add_textbox(slide, Inches(0.8), Inches(3.85), Inches(12), Inches(0.6),
                "Bons testes — e obrigado por validar o Facilitia.ia",
                size=22, color=EMERALD, bold=True)
    add_band(slide, Inches(0.8), Inches(4.6), Inches(2.0), Inches(0.06), EMERALD)
    add_textbox(slide, Inches(0.8), Inches(4.85), Inches(12), Inches(0.45),
                "Abril 2026 · Sprints 1 a 5 · 79 casos de uso · 100% aprovados · V3 comercial",
                size=14, color=TXT_MUTED)
    # Métricas finais
    metrics = [
        ("96", "Slides"),
        ("79", "Casos de uso"),
        ("322", "Testes"),
        ("217", "Regras"),
    ]
    x = Inches(0.8)
    for num, label in metrics:
        add_rect(slide, x, Inches(5.6), Inches(2.7), Inches(1.1),
                 fill=BG_CARD, border=BORDER_SOFT)
        add_textbox(slide, x, Inches(5.75), Inches(2.7), Inches(0.5),
                    num, size=28, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.25), Inches(2.7), Inches(0.35),
                    label, size=11, color=TXT_MUTED, align=PP_ALIGN.CENTER)
        x += Inches(2.9)
    add_footer(slide, "Facilitia.ia · Apresentação comercial de validação · Sprints 1-5")

    # ── Slide oculto com roteiro de narração ──
    slide = new_slide(prs)
    add_pill(slide, Inches(0.6), Inches(0.35), Inches(4.5), Inches(0.32),
             "📻  ROTEIRO DE NARRAÇÃO · NotebookLM",
             fill=EMERALD, text_color=WHITE, size=10, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.78), Inches(12.2), Inches(0.6),
                "Roteiro de narração", size=26, color=TXT_INK, bold=True)
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.35),
                "Texto para gerar áudio/vídeo de apresentação",
                size=13, color=TXT_MUTED)
    roteiro = """O Facilitia.ia é um sistema completo de gestão de licitações que cobre desde a descoberta de um edital até a execução do contrato e o CRM pós-venda. Foi desenvolvido em 5 sprints, cada uma construindo sobre a anterior.

PROCESSO DE VALIDAÇÃO: A automação Playwright executou todos os testes com a empresa CH Hospitalar (tutorial-1), gerando screenshots como prova. O validador humano recebe: os requisitos completos, os casos de uso, o roteiro tutorial-2 (empresa RP3X), e os resultados da validação automática.

SPRINT 1 — FUNDAÇÃO: Cadastro da empresa, portfolio de produtos (com IA que consulta ANVISA), e todas as configurações: pesos de score, custos, regiões, fontes de busca. Screenshots mostram dados reais da CH Hospitalar salvos com sucesso.

SPRINT 2 — CAPTAÇÃO: Busca automática no PNCP, scores em 6 dimensões (técnico, comercial, jurídico, logístico, documental, complexidade), decisão GO/NO-GO com IA. Screenshots mostram resultados de busca, scores e análises reais.

SPRINT 3 — PREÇO E PROPOSTA: 6 camadas de precificação (A=volumetria até F=histórico), motor de proposta técnica com IA, auditoria ANVISA, checklist documental, exportação de dossiê completo.

SPRINT 4 — DEFESA JURÍDICA: Análise legal do edital contra Lei 14.133/2021, tabela de inconsistências com gravidade, petição de impugnação gerada pela IA, controle de prazo com countdown, recursos e contra-razões.

SPRINT 5 — PÓS-VENDA: Follow-up (vitória/derrota), atas de pregão, gestão de contratos com empenhos (incluindo itens sem valor com alerta), auditoria cruzada, contratos a vencer em 5 tiers, CRM com pipeline de 13 stages, parametrizações, mapa por UF, agenda com urgências, KPIs, decisões de não-participação e perda.

O validador deve seguir os tutorialsprint#-2.md na ordem (Sprint 1 a 5), usando os dados fornecidos para a empresa RP3X (valida2@valida.com.br / 123456). Cada tela deve mostrar dados reais — se algo estiver vazio, reportar como problema."""
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(12.2), Inches(4.9),
             fill=BG_CARD, border=BORDER_SOFT)
    add_textbox(slide, Inches(0.85), Inches(2.2), Inches(11.8), Inches(4.6),
                roteiro, size=10, color=TXT_BODY)
    add_footer(slide, "Roteiro oculto · uso em NotebookLM ou narração comercial")

    # ── Salvar ──
    prs.save(OUT)
    print(f"✅ PowerPoint salvo: {OUT}")
    print(f"   Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()